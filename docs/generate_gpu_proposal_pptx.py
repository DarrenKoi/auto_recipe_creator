"""GUI VLM 벤치마크 비교 및 H200 GPU 도입 제안 PPTX 생성 스크립트

현황: Kimi-K2.5는 사내 서버에서 자체 운영 중 (API 제공)
제안: GUI 자동화 전용 모델(UI-TARS, MAI-UI, UI-Venus 등) + UI 파싱 도구(OmniParser 등)를
      별도 H200 GPU 2대에서 운영
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ──
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
SURFACE = RGBColor(0x1E, 0x29, 0x3B)
SURFACE2 = RGBColor(0x33, 0x41, 0x55)
WHITE = RGBColor(0xE2, 0xE8, 0xF0)
DIM = RGBColor(0x94, 0xA3, 0xB8)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
ACCENT2 = RGBColor(0x81, 0x8C, 0xF8)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
RED = RGBColor(0xF8, 0x71, 0x71)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)


def set_slide_bg(slide, color=BG_DARK):
    """슬라이드 배경색 설정"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Malgun Gothic"):
    """텍스트 박스 추가 헬퍼"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_shape_box(slide, left, top, width, height, fill_color=SURFACE, border_color=None):
    """배경 사각형 추가"""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_bullet_frame(slide, left, top, width, height, items, font_size=14, color=DIM):
    """글머리 기호 텍스트 프레임"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Malgun Gothic"
        p.space_after = Pt(4)
    return txBox


def add_table_row(slide, x_start, y, col_widths, values, colors, row_h=0.36,
                  bg=SURFACE, font_size=10, bold_cols=None):
    """테이블 행 추가"""
    add_shape_box(slide, x_start, y, sum(col_widths) + 0.4, row_h, bg)
    x = x_start + 0.1
    for j, (val, col_w) in enumerate(zip(values, col_widths)):
        is_bold = bold_cols and j in bold_cols
        add_textbox(slide, x, y, col_w, row_h, val,
                    font_size=font_size,
                    color=colors[j] if j < len(colors) else DIM,
                    bold=is_bold)
        x += col_w


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    row_h = 0.36

    # ════════════════════════════════════════
    # SLIDE 1: 표지
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 1.5, 1.2, 10, 1.2,
                "GUI 자동화 전용 VLM 및 UI 파싱 모델 도입 제안",
                font_size=34, color=ACCENT, bold=True)
    add_textbox(slide, 1.5, 2.5, 10, 0.8,
                "CD-SEM / RCS 레시피 자동화를 위한 GPU H200 x 2 도입 필요성",
                font_size=22, color=ACCENT2)
    add_textbox(slide, 1.5, 3.5, 10, 0.8,
                "벤치마크 기반 분석: Kimi-K2.5 (현재) vs GUI 전문 모델 + UI 파싱 도구",
                font_size=16, color=DIM)
    add_textbox(slide, 1.5, 4.5, 10, 0.8,
                "현황: Kimi-K2.5는 사내 서버에서 자체 운영 중 (API 서빙)\n"
                "제안: GUI 자동화 전용 모델 + OmniParser를 별도 GPU에서 운영",
                font_size=14, color=YELLOW)
    add_textbox(slide, 1.5, 6.0, 10, 0.5,
                "2026년 2월",
                font_size=14, color=DIM)

    # ════════════════════════════════════════
    # SLIDE 2: 문제점
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.3, 12, 0.8,
                "문제: Kimi-K2.5로는 복잡한 엔지니어링 UI를 정밀하게 처리할 수 없음",
                font_size=26, color=RED, bold=True)

    # 왼쪽 박스
    add_shape_box(slide, 0.8, 1.3, 5.8, 5.5, SURFACE, SURFACE2)
    add_textbox(slide, 1.0, 1.4, 5.4, 0.5,
                "Kimi-K2.5 현재 성능 한계", font_size=18, color=RED, bold=True)
    items = [
        "  ScreenSpot-Pro 벤치마크 정확도: 52.8%",
        "  GUI 전용 SOTA 모델 (69.6%) 대비 17% 부족",
        "  범용 모델 — GUI 인터랙션 데이터로 학습되지 않음",
        "  아이콘/위젯 감지 능력이 엔지니어링 SW에서 취약",
        "  사내 운영 중이나, 범용 모델의 근본적 한계 존재",
        "  RCS/CD-SEM 전용 파인튜닝이 어려운 구조",
        "",
        "  * OSWorld (63.3%)에서는 높은 점수이나,",
        "    이는 다단계 작업 계획 능력이지",
        "    정밀 클릭 좌표 능력이 아님",
    ]
    add_bullet_frame(slide, 1.0, 2.0, 5.4, 4.5, items, font_size=13, color=DIM)

    # 오른쪽 박스
    add_shape_box(slide, 6.8, 1.3, 5.8, 5.5, SURFACE, RED)
    add_textbox(slide, 7.0, 1.4, 5.4, 0.5,
                "핵심 문제: 좁은 영역 UI 요소 인식 실패", font_size=18, color=RED, bold=True)
    items2 = [
        "  UI 요소가 좁은 영역에 밀집되어 있을 때 (<20px 간격),",
        "  Kimi-K2.5는 인접 요소를 구분하지 못합니다:",
        "",
        "  - 클릭 좌표가 인접 버튼/라벨로 이탈",
        "  - 툴바 아이콘 클러스터에서 모호한 좌표 출력",
        "  - TreeView / ListView 행 높이가 작으면 행 혼동",
        "  - RCS 레시피 편집기는 <10px 간격의 밀집 패널 보유",
        "",
        "  근본 원인: GUI 공간 학습 부재",
        "  범용 VLM은 자연 이미지/문서로 학습됨",
        "  → 밀집 엔지니어링 UI에서 타겟 크기 줄어들수록",
        "    정확도가 급격히 하락",
    ]
    add_bullet_frame(slide, 7.0, 2.0, 5.4, 4.5, items2, font_size=12, color=DIM)

    # ════════════════════════════════════════
    # SLIDE 3: ScreenSpot-Pro 벤치마크
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.3, 11, 0.8,
                "ScreenSpot-Pro: 전문 GUI 그라운딩 벤치마크",
                font_size=28, color=ACCENT, bold=True)
    add_textbox(slide, 0.8, 1.0, 11, 0.5,
                "고해상도 전문 애플리케이션 UI (CAD, EDA, 개발 도구)에서의 클릭 좌표 정확도 — "
                "CD-SEM/RCS 유즈케이스에 가장 유사한 공개 벤치마크",
                font_size=13, color=DIM)

    models = [
        ("UI-Venus-1.5-8B (Huawei)", "GUI 전용", "8B", "69.6", GREEN),
        ("MAI-UI-32B (Alibaba)", "GUI 전용", "32B", "67.9", GREEN),
        ("MAI-UI-8B (Alibaba)", "GUI 전용", "8B", "65.7", GREEN),
        ("Qwen3-VL-30B-A3B", "범용", "30B MoE", "61.8", ACCENT),
        ("UI-TARS-1.5-7B (ByteDance)", "GUI 전용", "7B", "61.6", ACCENT),
        ("MAI-UI-2B (Alibaba)", "GUI 전용", "2B", "57.4", YELLOW),
        ("Qwen3-VL-8B", "범용", "8B", "54.6", YELLOW),
        ("Kimi-K2.5 (현재 사내 운영 중)", "범용", "~200B+ MoE", "52.8", ORANGE),
        ("GUI-Actor-7B (Microsoft)", "GUI 전용", "7B", "44.6", ORANGE),
        ("OmniParser V2 + GPT-4o", "파서+VLM", "YOLO+Florence", "39.6", TEAL),
        ("UI-TARS-72B v1", "GUI 전용", "72B", "38.1", ORANGE),
        ("GPT-4o (단독)", "범용", "~200B+", "0.8", RED),
    ]

    headers = ["모델", "유형", "파라미터", "ScreenSpot-Pro (%)"]
    col_widths = [4.2, 1.3, 1.8, 2.2]
    x_start = 0.8
    y_start = 1.6

    add_shape_box(slide, x_start, y_start, sum(col_widths) + 0.4, row_h, SURFACE2)
    x = x_start + 0.1
    for header, w in zip(headers, col_widths):
        add_textbox(slide, x, y_start, w, row_h, header,
                    font_size=11, color=ACCENT, bold=True)
        x += w

    for i, (model, mtype, params, score, color) in enumerate(models):
        y = y_start + row_h * (i + 1)
        bg = SURFACE if i % 2 == 0 else BG_DARK
        if "현재" in model:
            bg = RGBColor(0x3B, 0x2A, 0x1A)
        add_shape_box(slide, x_start, y, sum(col_widths) + 0.4, row_h, bg)

        x = x_start + 0.1
        is_current = "현재" in model
        add_textbox(slide, x, y, col_widths[0], row_h, model,
                    font_size=10, color=ORANGE if is_current else WHITE,
                    bold=is_current)
        x += col_widths[0]
        type_color = GREEN if "GUI" in mtype else (TEAL if "파서" in mtype else YELLOW)
        add_textbox(slide, x, y, col_widths[1], row_h, mtype, font_size=10, color=type_color)
        x += col_widths[1]
        add_textbox(slide, x, y, col_widths[2], row_h, params, font_size=10, color=DIM)
        x += col_widths[2]
        add_textbox(slide, x, y, col_widths[3], row_h, score + "%",
                    font_size=12, color=color, bold=True)

    # 인사이트 박스
    add_shape_box(slide, 10.2, 1.6, 2.8, 5.5, SURFACE, ACCENT)
    add_textbox(slide, 10.3, 1.7, 2.6, 0.4,
                "핵심 수치", font_size=14, color=ACCENT, bold=True)
    insight_items = [
        "현재 Kimi-K2.5: 52.8%",
        "",
        "GUI 전용 SOTA:",
        "UI-Venus-1.5: 69.6%",
        "→ +32% 정확도 향상",
        "",
        "UI-TARS-1.5-7B: 61.6%",
        "→ +17% 향상 (7B만으로)",
        "",
        "OmniParser가 GPT-4o를",
        "0.8% → 39.6%로 끌어올림",
        "→ 파서+VLM 조합의 힘",
    ]
    add_bullet_frame(slide, 10.3, 2.2, 2.6, 4.5, insight_items, font_size=11, color=DIM)

    # ════════════════════════════════════════
    # SLIDE 4: UI 파싱 도구 + 파이프라인
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.3, 11, 0.8,
                "해결 전략: UI 파싱 도구 + GUI 전용 VLM 결합",
                font_size=28, color=TEAL, bold=True)
    add_textbox(slide, 0.8, 1.0, 11, 0.5,
                "UI 파싱 도구가 스크린샷의 모든 인터랙션 요소를 사전 감지 → GUI 전용 VLM이 정확한 좌표 결정",
                font_size=13, color=DIM)

    # 파이프라인 다이어그램
    add_shape_box(slide, 0.8, 1.6, 11.7, 1.2, SURFACE2, TEAL)
    add_textbox(slide, 1.0, 1.65, 11.3, 0.4,
                "제안 파이프라인 (2단계 아키텍처)", font_size=14, color=TEAL, bold=True)
    add_textbox(slide, 1.0, 2.05, 11.3, 0.6,
                "스크린샷 → [OmniParser V2: YOLOv8 + Florence-2 + OCR] → 구조화된 요소 (bbox + 라벨) "
                "→ [GUI 전용 VLM: UI-TARS / UI-Venus / MAI-UI] → 정확한 클릭 좌표 + 액션 → [pynput 마우스/키보드]",
                font_size=12, color=DIM)

    # OmniParser 상세
    add_shape_box(slide, 0.8, 3.0, 5.8, 3.8, SURFACE, TEAL)
    add_textbox(slide, 1.0, 3.05, 5.4, 0.4,
                "OmniParser V2 (Microsoft, 오픈소스)", font_size=16, color=TEAL, bold=True)
    omni_items = [
        "아키텍처:",
        "  1단계: YOLOv8 Nano — 인터랙션 요소 감지 (bbox)",
        "  2단계: Florence-2 — 아이콘 기능 설명 생성",
        "  3단계: PaddleOCR — 텍스트 추출",
        "",
        "성능:",
        "  - GPT-4o 단독: 0.8% → OmniParser 결합: 39.6%",
        "  - V1 대비 지연시간 60% 감소 (A100 기준 0.6초/프레임)",
        "  - 작은 요소 (<20px) 감지 정확도 대폭 향상",
        "",
        "OmniTool: Docker 기반 Windows 11 VM 제공",
        "  → 에이전트 테스트 + 학습 데이터 자동 수집 가능",
    ]
    add_bullet_frame(slide, 1.0, 3.5, 5.4, 3.2, omni_items, font_size=11, color=DIM)

    # 기타 파싱 도구
    add_shape_box(slide, 6.8, 3.0, 5.8, 3.8, SURFACE, ACCENT)
    add_textbox(slide, 7.0, 3.05, 5.4, 0.4,
                "기타 UI 파싱 도구 및 기법", font_size=16, color=ACCENT, bold=True)
    other_items = [
        "SparkUI-Parser (Ant Group, 오픈소스):",
        "  - MLLM + 토큰 라우터 + 좌표 디코더",
        "  - 연속 좌표 모델링 (이산화 없음) → 최고 정밀도",
        "  - 클래스 최고 추론 속도",
        "",
        "Set-of-Mark / SoM (Microsoft):",
        "  - 감지된 UI 요소에 번호 마크 오버레이",
        "  - VLM이 좌표 대신 번호로 요소 참조",
        "  - 모든 감지 모델과 조합 가능",
        "",
        "SeeClick (NJU):",
        "  - Qwen-VL 기반 파인튜닝",
        "  - UI 스크린샷만으로 자동 GUI 인식의 첫 사례",
    ]
    add_bullet_frame(slide, 7.0, 3.5, 5.4, 3.2, other_items, font_size=11, color=DIM)

    # ════════════════════════════════════════
    # SLIDE 5: Agent 벤치마크
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.3, 11, 0.8,
                "End-to-End Agent 벤치마크: 다단계 데스크톱 자동화",
                font_size=28, color=ACCENT, bold=True)
    add_textbox(slide, 0.8, 1.0, 11, 0.5,
                "OSWorld & WindowsAgentArena = Windows 데스크톱 자동화 직접 테스트 — RCS/CD-SEM과 동일한 환경",
                font_size=13, color=DIM)

    agent_models = [
        ("Kimi-K2.5 (현재 사내 운영)", "63.3", "—", "—", "—"),
        ("UI-Venus-1.5 (Huawei)", "—", "—", "77.6", "—"),
        ("UI-TARS-2 (ByteDance)", "47.5", "50.6", "73.3", "88.2"),
        ("Claude 4 Sonnet", "43.9", "—", "—", "—"),
        ("OpenAI CUA-o3", "42.9", "—", "52.5", "71.0"),
        ("UI-TARS-1.5 (ByteDance)", "42.5", "42.1", "64.2", "75.8"),
        ("Qwen3-VL-32B", "—", "—", "63.7", "—"),
    ]

    headers2 = ["모델", "OSWorld", "WinAgentArena", "AndroidWorld", "Mind2Web"]
    col_widths2 = [3.8, 1.6, 2.0, 2.0, 1.8]
    x_start2 = 1.0
    y_start2 = 1.7

    add_shape_box(slide, x_start2, y_start2, sum(col_widths2) + 0.4, row_h, SURFACE2)
    x = x_start2 + 0.1
    for header, w in zip(headers2, col_widths2):
        add_textbox(slide, x, y_start2, w, row_h, header,
                    font_size=12, color=ACCENT, bold=True)
        x += w

    for i, (model, *scores) in enumerate(agent_models):
        y = y_start2 + row_h * (i + 1)
        bg = SURFACE if i % 2 == 0 else BG_DARK
        if "현재" in model:
            bg = RGBColor(0x3B, 0x2A, 0x1A)
        add_shape_box(slide, x_start2, y, sum(col_widths2) + 0.4, row_h, bg)
        x = x_start2 + 0.1
        is_current = "현재" in model
        add_textbox(slide, x, y, col_widths2[0], row_h, model,
                    font_size=11, color=ORANGE if is_current else WHITE, bold=True)
        x += col_widths2[0]
        for j, s in enumerate(scores):
            c = GREEN if s not in ("—", "") and float(s) > 45 else (ACCENT if s not in ("—", "") else DIM)
            add_textbox(slide, x, y, col_widths2[j + 1], row_h, s,
                        font_size=12, color=c, bold=(s not in ("—", "")))
            x += col_widths2[j + 1]

    # 하단 설명
    add_shape_box(slide, 1.0, 4.6, 11.3, 2.3, SURFACE, ACCENT)
    add_textbox(slide, 1.2, 4.65, 10.9, 0.4,
                "Kimi-K2.5 OSWorld 점수가 높은데 왜 부족한가?", font_size=16, color=YELLOW, bold=True)
    note_items = [
        "  Kimi-K2.5는 OSWorld (63.3%)에서 높은 점수 → 다단계 작업 계획/추론 능력은 우수",
        "  그러나 ScreenSpot-Pro (52.8%)에서 드러나듯, 정밀 클릭 좌표 능력은 GUI 전용 모델 대비 크게 부족",
        "  RCS 자동화의 병목은 '어떤 순서로 할 것인가' (계획)가 아니라 '정확히 어디를 클릭할 것인가' (그라운딩)",
        "  GUI 전용 모델 (UI-TARS, UI-Venus)은 정밀 그라운딩에 특화 → Kimi-K2.5의 계획 능력과 상호 보완 가능",
    ]
    add_bullet_frame(slide, 1.2, 5.1, 10.9, 1.8, note_items, font_size=12, color=DIM)

    # ════════════════════════════════════════
    # SLIDE 6: 왜 H200 x 2
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.3, 11, 0.8,
                "GPU 리소스 요구사항: H200 x 2 도입 근거",
                font_size=28, color=ACCENT, bold=True)
    add_textbox(slide, 0.8, 0.95, 11, 0.4,
                "기존 Kimi-K2.5 사내 서버와 별도로, GUI 자동화 전용 모델 + UI 파싱 도구를 위한 독립 GPU 인프라",
                font_size=13, color=YELLOW)

    # GPU 스펙 + 배포 전략
    add_shape_box(slide, 0.8, 1.5, 6.0, 2.7, SURFACE, ACCENT)
    add_textbox(slide, 1.0, 1.55, 5.6, 0.4,
                "NVIDIA H200 SXM 사양 및 배포 전략", font_size=16, color=ACCENT, bold=True)
    specs = [
        "  HBM3e 메모리: GPU당 141 GB / 2대 합계 282 GB",
        "  메모리 대역폭: 4.8 TB/s | FP16: 989 TFLOPS",
        "",
        "  GPU 1 배포 계획:",
        "    OmniParser V2 (~4GB) + UI-TARS-1.5-7B (~14GB)",
        "    + UI-Venus-1.5-8B (~16GB) = 약 34GB → 배치 추론 여유",
        "",
        "  GPU 2 배포 계획:",
        "    MAI-UI-32B (~64GB) 또는 파인튜닝 워크로드",
        "    또는 UI-TARS-72B (~72GB INT8) 대형 모델 운영",
    ]
    add_bullet_frame(slide, 1.0, 2.0, 5.6, 2.1, specs, font_size=11, color=DIM)

    # 2대 근거
    add_shape_box(slide, 7.0, 1.5, 5.8, 2.7, SURFACE, GREEN)
    add_textbox(slide, 7.2, 1.55, 5.4, 0.4,
                "왜 2대인가 (1대가 아닌 이유)", font_size=16, color=GREEN, bold=True)
    reasons = [
        "1. 파서 + 다중 VLM 동시 운영:",
        "   OmniParser + 2~3개 GUI 전용 모델 병렬 서빙",
        "2. 대형 모델 접근: 32B/72B 모델은 64~144GB 필요",
        "   → 1대로는 다른 모델과 동시 운영 불가",
        "3. A/B 테스트: UI-TARS vs MAI-UI vs UI-Venus",
        "   비교 평가 및 자동 폴백 구현",
        "4. 파인튜닝: LoRA 학습 시 추론의 2~4배 메모리",
        "   → 한쪽에서 서빙, 한쪽에서 학습 병행",
        "5. 미래 대비: 차세대 모델은 더 큰 파라미터 추세",
    ]
    add_bullet_frame(slide, 7.2, 2.0, 5.4, 2.1, reasons, font_size=11, color=DIM)

    # VRAM 테이블
    vram_data = [
        ("OmniParser V2 (파서)", "YOLO+Florence", "~4-8 GB", "~2-4 GB"),
        ("UI-Venus-1.5-8B", "8B", "~16 GB", "~8 GB"),
        ("UI-TARS-1.5-7B", "7B", "~14 GB", "~7 GB"),
        ("MAI-UI-8B", "8B", "~16 GB", "~8 GB"),
        ("MAI-UI-32B (SOTA)", "32B", "~64 GB", "~32 GB"),
        ("UI-TARS-1.5-72B", "72B", "~144 GB", "~72 GB"),
    ]

    y_table = 4.5
    headers3 = ["모델", "파라미터", "FP16 VRAM", "INT8 VRAM"]
    col_w3 = [3.5, 1.8, 1.8, 1.8]
    add_shape_box(slide, 0.8, y_table, sum(col_w3) + 0.4, 0.35, SURFACE2)
    x = 0.9
    for h, w in zip(headers3, col_w3):
        add_textbox(slide, x, y_table, w, 0.35, h, font_size=11, color=ACCENT, bold=True)
        x += w

    for i, (model, params, fp16, int8) in enumerate(vram_data):
        y = y_table + 0.35 * (i + 1)
        bg = SURFACE if i % 2 == 0 else BG_DARK
        add_shape_box(slide, 0.8, y, sum(col_w3) + 0.4, 0.35, bg)
        x = 0.9
        add_textbox(slide, x, y, col_w3[0], 0.35, model, font_size=10, color=WHITE)
        x += col_w3[0]
        add_textbox(slide, x, y, col_w3[1], 0.35, params, font_size=10, color=DIM)
        x += col_w3[1]
        add_textbox(slide, x, y, col_w3[2], 0.35, fp16, font_size=10, color=DIM)
        x += col_w3[2]
        add_textbox(slide, x, y, col_w3[3], 0.35, int8, font_size=10, color=DIM)

    # ════════════════════════════════════════
    # SLIDE 7: 현행 vs 제안
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.3, 12, 0.8,
                "현행 대비 개선 효과",
                font_size=28, color=ACCENT, bold=True)

    # 왼쪽: 현행
    add_shape_box(slide, 0.8, 1.2, 5.8, 5.8, SURFACE, RED)
    add_textbox(slide, 1.0, 1.3, 5.4, 0.4,
                "현행: Kimi-K2.5 단독 (사내 서버)", font_size=18, color=RED, bold=True)
    api_items = [
        "운영 환경:",
        "  - 사내 서버에서 Kimi-K2.5 API 서빙 중",
        "  - 범용 VLM으로 다양한 업무에 활용",
        "",
        "GUI 자동화 한계:",
        "  - ScreenSpot-Pro 정확도 52.8%",
        "  - 밀집된 좁은 영역 UI 요소 인식 불가",
        "  - 클릭 좌표가 인접 요소로 빈번하게 이탈",
        "  - UI 파싱 전처리 없이 VLM에 직접 의존",
        "  - 범용 모델의 근본적 한계 — 해결 불가",
        "",
        "결론:",
        "  - Kimi-K2.5는 범용 업무 및 작업 계획에 적합",
        "  - 정밀 GUI 그라운딩에는 전용 모델이 필요",
        "  - UI 파싱 도구로 사전 요소 감지가 필수",
    ]
    add_bullet_frame(slide, 1.0, 1.8, 5.4, 5.0, api_items, font_size=12, color=DIM)

    # 오른쪽: 제안
    add_shape_box(slide, 6.8, 1.2, 5.8, 5.8, SURFACE, GREEN)
    add_textbox(slide, 7.0, 1.3, 5.4, 0.4,
                "제안: 파서 + GUI VLM (별도 H200 x 2)", font_size=18, color=GREEN, bold=True)
    self_items = [
        "운영 구조:",
        "  - 기존 Kimi-K2.5 서버는 그대로 유지",
        "  - 별도 H200 GPU 2대에 전용 모델 배포",
        "  - OmniParser (파서) + GUI VLM 2단계 파이프라인",
        "",
        "기대 효과:",
        "  - 정확도 61.6~69.6% (UI-TARS, UI-Venus)",
        "  - 현재 대비 최대 +32% 향상",
        "  - OmniParser가 밀집 UI 요소를 사전 감지",
        "  - 좁은 영역도 bbox로 분리 → 혼동 해소",
        "  - RCS/CD-SEM 스크린샷으로 파인튜닝 가능",
        "",
        "추가 이점:",
        "  - Kimi-K2.5와 독립 운영 → 기존 서비스 무영향",
        "  - OmniTool로 학습 데이터 자동 수집 가능",
        "  - 멀티 모델 앙상블/폴백으로 안정성 확보",
    ]
    add_bullet_frame(slide, 7.0, 1.8, 5.4, 5.0, self_items, font_size=12, color=DIM)

    # ════════════════════════════════════════
    # SLIDE 8: 구현 로드맵
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 0.8, 0.3, 11, 0.8,
                "구현 로드맵",
                font_size=28, color=ACCENT, bold=True)

    phases = [
        ("1단계: 즉시 착수 (1~2주차)", ACCENT,
         ["H200 GPU 1에 OmniParser V2 + UI-TARS-1.5-7B 배포 (vLLM 서빙)",
          "OmniParser → UI-TARS 2단계 파이프라인 구현, 기존 poc/work/ 코드에 통합",
          "A/B 테스트: Kimi-K2.5 단독 vs OmniParser+UI-TARS (RCS 로그인 + 탭 전환)"]),
        ("2단계: 검증 (3~4주차)", YELLOW,
         ["GPU 2에 UI-Venus-1.5-8B / MAI-UI-8B 배포 (보조 모델)",
          "실제 RCS 스크린샷으로 벤치마크: 밀집 컨트롤 패널, 좁은 영역 클릭 정확도",
          "모델 앙상블/폴백 구현: 파서+UI-TARS 우선 → UI-Venus 보조"]),
        ("3단계: 최적화 (2~3개월차)", GREEN,
         ["자사 RCS/CD-SEM 스크린샷 데이터셋 구축 (OmniTool로 자동 수집)",
          "UI-TARS-1.5-7B 도메인 파인튜닝 (LoRA, GPU 2에서 학습)",
          "전체 RCS 레시피 생성 파이프라인에 통합"]),
        ("4단계: 프로덕션 (3개월차~)", ACCENT2,
         ["모니터링 및 자동 복구 기능 포함 프로덕션 배포",
          "차세대 모델 (UI-TARS-2, UI-Venus v2, MAI-UI v2) 출시 시 평가 및 교체",
          "내부 벤치마크 데이터셋으로 지속적 모델 평가 체계 확립"]),
    ]

    y = 1.2
    for title, color, items in phases:
        add_shape_box(slide, 0.8, y, 11.7, 0.35 + len(items) * 0.30, SURFACE, color)
        add_textbox(slide, 1.0, y + 0.02, 11, 0.35, title, font_size=14, color=color, bold=True)
        for j, item in enumerate(items):
            add_textbox(slide, 1.3, y + 0.35 + j * 0.28, 11, 0.28,
                        "  " + item, font_size=11, color=DIM)
        y += 0.45 + len(items) * 0.30

    # ════════════════════════════════════════
    # SLIDE 9: 요약 및 요청
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, 1.5, 0.6, 10, 1.0,
                "요약 및 요청사항",
                font_size=36, color=ACCENT, bold=True)

    summary = [
        "1.  현재 사내 운영 중인 Kimi-K2.5는 ScreenSpot-Pro 기준 52.8% 정확도",
        "     — 밀집 UI 요소가 많은 RCS 레시피 편집기에서 신뢰할 수 없는 수준",
        "",
        "2.  특히 좁은 영역에 밀집된 UI 요소 (<20px 간격)를 구분하지 못함",
        "     — 클릭 좌표가 인접 요소로 이탈하여 자동화 실패",
        "",
        "3.  GUI 전용 오픈소스 모델은 61~70% 정확도 (현재 대비 최대 +32% 향상)",
        "     — UI-TARS-1.5-7B (61.6%), UI-Venus-1.5-8B (69.6%, SOTA)",
        "",
        "4.  OmniParser 등 UI 파싱 도구를 전처리로 결합하면",
        "     밀집 요소 사전 감지 → VLM 그라운딩 정확도 추가 향상",
        "",
        "5.  기존 Kimi-K2.5 서버와 독립적으로 운영하여 기존 서비스 무영향",
        "     H200 GPU 2대 (282GB)로 파서 + 다중 모델 서빙 + 파인튜닝 가능",
    ]
    add_bullet_frame(slide, 1.5, 1.7, 10, 4.0, summary, font_size=14, color=DIM)

    add_shape_box(slide, 1.5, 5.6, 10.3, 1.2, SURFACE, ACCENT)
    add_textbox(slide, 1.7, 5.65, 9.9, 1.1,
                "요청: GUI 자동화 전용 VLM + UI 파싱 도구 배포 및 파인튜닝을 위한\n"
                "NVIDIA H200 SXM GPU 2대 도입을 승인해 주시기 바랍니다.\n"
                "(기존 Kimi-K2.5 인프라와 별도 운영, 기존 서비스에 영향 없음)",
                font_size=16, color=ACCENT, bold=True)

    # Save
    output_path = "/Users/daeyoung/Codes/auto_recipe_creator/docs/gui_vlm_gpu_h200_proposal.pptx"
    prs.save(output_path)
    print(f"[INFO] PPTX 생성 완료: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
