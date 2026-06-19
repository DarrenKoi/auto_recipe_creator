"""프로젝트 진행 보고서 PPT(.pptx) 생성 스크립트.

docs/project_progress/ 의 보고서 내용을 임원 보고용 슬라이드(~12장)로 작성한다.
시각 테마/헬퍼는 poc/workflow_1/build_report_pptx.py 의 NAVY/ACCENT 코퍼레이트 패턴을 재사용
(한글 폰트는 맑은 고딕으로 지정).

실행:
  uv run python docs/project_progress/build_report_pptx.py
출력:
  docs/project_progress/project_progress_report.pptx

수치/경로는 저장소 근거 문서(docs/setup_vlms, poc/workflow_2/docs, poc/workflow_3)에 따른다.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUTPUT_PATH = Path(__file__).resolve().parent / "project_progress_report.pptx"

FONT = "맑은 고딕"
NAVY = RGBColor(0x14, 0x2A, 0x55)
ACCENT = RGBColor(0xE8, 0x6A, 0x1F)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
TEXT = RGBColor(0x22, 0x29, 0x33)
MUTED = RGBColor(0x55, 0x60, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x1E, 0x7A, 0x3C)


def _set_run(run, *, size, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT


def _add_title_band(slide, title_text, subtitle_text):
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    band.shadow.inherit = False

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.13), Inches(12.5), Inches(0.55))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = title_text
    _set_run(run, size=25, bold=True, color=WHITE)

    sub_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.7), Inches(12.5), Inches(0.35))
    run = sub_box.text_frame.paragraphs[0].add_run()
    run.text = subtitle_text
    _set_run(run, size=12.5, color=RGBColor(0xCD, 0xD7, 0xE6))


def _add_section_card(slide, *, left, top, width, height, heading, bullets, accent=NAVY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = accent
    card.line.width = Pt(0.75)
    card.shadow.inherit = False

    head_box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.13), width - Inches(0.5), Inches(0.45))
    run = head_box.text_frame.paragraphs[0].add_run()
    run.text = heading
    _set_run(run, size=15, bold=True, color=accent)

    body_box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.62), width - Inches(0.5), height - Inches(0.75))
    tf = body_box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"• {bullet}"
        _set_run(run, size=11.5, color=TEXT)


def _add_pipeline_step(slide, *, left, top, width, height, index, title, detail):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.25)
    box.shadow.inherit = False

    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.15), top + Inches(0.15), Inches(0.4), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    badge.shadow.inherit = False
    btf = badge.text_frame
    btf.margin_left = btf.margin_right = btf.margin_top = btf.margin_bottom = 0
    p = btf.paragraphs[0]
    p.alignment = 2
    run = p.add_run()
    run.text = str(index)
    _set_run(run, size=12, bold=True, color=WHITE)

    title_box = slide.shapes.add_textbox(left + Inches(0.65), top + Inches(0.13), width - Inches(0.8), Inches(0.4))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = title
    _set_run(run, size=12.5, bold=True, color=NAVY)

    body_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), height - Inches(0.7))
    tf = body_box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = detail
    _set_run(run, size=10, color=MUTED)


def _add_loop_node(slide, *, left, top, width, height, title, detail, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    box.shadow.inherit = False

    title_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.35))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = title
    _set_run(run, size=11.5, bold=True, color=color)

    body_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.48), width - Inches(0.3), height - Inches(0.58))
    tf = body_box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = detail
    _set_run(run, size=9.5, color=MUTED)


def _add_arrow(slide, *, left, top, width, height, color, direction="right"):
    shape_map = {"right": MSO_SHAPE.RIGHT_ARROW, "down": MSO_SHAPE.DOWN_ARROW}
    arrow = slide.shapes.add_shape(shape_map[direction], left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    arrow.shadow.inherit = False


def _add_role_card(slide, *, left, top, width, height, role, model, responsibility, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    box.shadow.inherit = False

    role_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.35))
    run = role_box.text_frame.paragraphs[0].add_run()
    run.text = role
    _set_run(run, size=10.5, bold=True, color=color)

    model_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.44), width - Inches(0.4), Inches(0.35))
    run = model_box.text_frame.paragraphs[0].add_run()
    run.text = model
    _set_run(run, size=13, bold=True, color=NAVY)

    body_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.83), width - Inches(0.4), height - Inches(0.95))
    tf = body_box.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text = responsibility
    _set_run(run, size=10, color=TEXT)


def _add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    _set_run(run, size=10, color=MUTED)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# 슬라이드
# ---------------------------------------------------------------------------


def slide_cover(prs):
    slide = _blank(prs)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.4), Inches(13.333), Inches(2.7))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    band.shadow.inherit = False

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.9))
    run = title_box.text_frame.paragraphs[0].add_run()
    run.text = "프로젝트 진행 보고서"
    _set_run(run, size=40, bold=True, color=WHITE)

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.6))
    run = sub_box.text_frame.paragraphs[0].add_run()
    run.text = "AI 기반 CD-SEM / VeritySEM Recipe 자동 Setup PoC"
    _set_run(run, size=18, color=RGBColor(0xCD, 0xD7, 0xE6))

    foot_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5))
    run = foot_box.text_frame.paragraphs[0].add_run()
    run.text = "VLM 배포·운영  ·  workflow_1 / workflow_2 / workflow_3      |      목적 · PoC 방향 · 성과 · 확장성"
    _set_run(run, size=13, color=ACCENT)


def slide_problem_direction(prs):
    slide = _blank(prs)
    _add_title_band(slide, "문제 정의 & PoC 핵심 방향",
                    "수동 recipe setup의 한계 → VLM(이해)과 CV(좌표)의 역할 분리")
    _add_section_card(slide, left=Inches(0.4), top=Inches(1.4), width=Inches(6.25), height=Inches(2.4),
                      heading="문제 정의", bullets=[
                          "CD-SEM recipe setup을 사람이 수동으로 — 야간/주말 무인 불가",
                          "RCS는 legacy GUI(UIA 부실) → 일반 자동화로 안정 제어 어려움",
                          "Align key는 contrast·밝기·반복패턴이 매 측정마다 달라 단순 인식 불가",
                          "실패 원인 분석용 화면 증거도 사후 소실",
                      ])
    _add_section_card(slide, left=Inches(6.85), top=Inches(1.4), width=Inches(6.1), height=Inches(2.4),
                      heading="PoC 방향 — 역할 분리", accent=ACCENT, bullets=[
                          "VLM = '어디를 볼까' (영역 식별·모호성 설명·모드 판독)",
                          "CV(OpenCV) = '얼마나 닮았나 / 정확히 어디' (정량 점수·최종 좌표)",
                          "낮은 CV 점수를 VLM이 덮어쓰지 않음",
                          "stateless VLM의 '기억'을 CV 점수가 외부 상태로 대신",
                      ])
    _add_section_card(slide, left=Inches(0.4), top=Inches(3.95), width=Inches(12.55), height=Inches(2.9),
                      heading="왜 역할 분리가 핵심인가", bullets=[
                          "전체 화면에서 클릭 좌표를 한 번에 픽셀 정확도로 맞히기 어렵다 → coarse(VLM) + fine(VLM) + 정량검증(CV)",
                          "VLM은 호출 간 기억이 없다(stateless) → '직전 위치/유사도'를 CV 점수로 외부 보존해 다음 단계로 전달",
                          "화면 이해(VLM)와 정밀 좌표(CV)를 각자 강점에 맡겨 무인 자동화 신뢰도 확보",
                      ])
    _add_footer(slide, "설계 원칙 확정 2026-05-25  •  VLM은 식별/판단, CV는 점수/좌표")


def slide_timeline(prs):
    slide = _blank(prs)
    _add_title_band(slide, "전체 진행 흐름 (Workstream Timeline)",
                    "VLM 인프라(0) → GUI 자동화 증명(1) → CV 정확도(2) → 실시간 통합 루프(3)")
    steps = [
        ("0. deploy_vlms", "오픈소스 VLM 5종 조사·선정 → 사내 HCP(H200×2)에 vLLM 설치, Flask proxy 통합 운영"),
        ("1. workflow_1", "RCS GUI 자동화(2-stage VLM) + Align Fail 감지 + CCTV 캡처 PoC (검증 후 동결)"),
        ("2. workflow_2", "오프라인 CV 평가 벤치 — golden set으로 matching/consensus A/B·튜닝"),
        ("3. workflow_3", "통합 production 실시간 align-fail 모니터링 루프 (현재 주력)"),
    ]
    top = Inches(2.0)
    height = Inches(2.6)
    width = Inches(2.9)
    gap = Inches(0.25)
    arrow_w = Inches(0.22)
    left = Inches(0.5)
    for idx, (title, detail) in enumerate(steps):
        _add_loop_node(slide, left=left, top=top, width=width, height=height, title=title, detail=detail,
                       color=NAVY if idx % 2 == 0 else ACCENT)
        if idx < len(steps) - 1:
            _add_arrow(slide, left=left + width + Inches(0.015), top=top + (height - Inches(0.45)) / 2,
                       width=arrow_w, height=Inches(0.45), color=NAVY)
        left += width + gap
    _add_section_card(slide, left=Inches(0.5), top=Inches(5.1), width=Inches(12.3), height=Inches(1.75),
                      heading="흐름 요약", accent=ACCENT, bullets=[
                          "각 단계가 다음 단계의 전제를 만든다: 인프라 → 자동화 가능성 → 정확도 → 실시간 통합",
                          "검증된 CV 변경만 wf2(벤치)에서 wf3(production)로 bit-parity 포팅 → 회귀 위험 통제",
                      ])
    _add_footer(slide, "deploy_vlms → workflow_1 → workflow_2 → workflow_3")


def slide_vlm_roles(prs):
    slide = _blank(prs)
    _add_title_band(slide, "VLM 배포·운영 — 5개 특화 모델 분담",
                    "단일 거대 모델 대신 역할별 오픈소스 VLM을 Flask proxy로 통합")
    roles = [
        ("Port 8001", "UI-Venus-1.5-8B", "메인 GUI grounding — 전체 화면 coarse bbox", NAVY),
        ("Port 8002", "MAI-UI-8B", "정밀 클릭점 — coarse crop 확대 후 픽셀 좌표", ACCENT),
        ("Port 8004", "PaddleOCR-VL-1.5", "OCR 보조 — 텍스트·Spotting(좌표)·표 인식", NAVY),
        ("Port 8005 / 8003", "GOT-OCR / UI-TARS", "OCR fallback / GUI agent 대안 경로", ACCENT),
    ]
    top = Inches(1.45)
    height = Inches(2.2)
    width = Inches(3.05)
    gap = Inches(0.13)
    left = Inches(0.4)
    for role, model, resp, color in roles:
        _add_role_card(slide, left=left, top=top, width=width, height=height, role=role, model=model,
                       responsibility=resp, color=color)
        left += width + gap
    _add_section_card(slide, left=Inches(0.4), top=Inches(3.85), width=Inches(6.25), height=Inches(3.0),
                      heading="사내 HCP 설치", bullets=[
                          "하드웨어: H200 140 GiB × 2, 서빙: vLLM(BF16) + GOT-OCR(transformers)",
                          "GPU 0: UI-Venus + UI-TARS(auto-tune u≈0.44) → ~123 GiB(88%)",
                          "GPU 1: MAI-UI + PaddleOCR-VL + GOT-OCR → ~81 GiB(58%)",
                          "오프라인 정책(HF live pull 금지·telemetry off), prefix caching 활용",
                      ])
    _add_section_card(slide, left=Inches(6.85), top=Inches(3.85), width=Inches(6.1), height=Inches(3.0),
                      heading="Flask Proxy 통합", accent=ACCENT, bullets=[
                          "서비스 레지스트리(config.py): route_slug / port / enabled",
                          "URL: {base}/api/vlm_serve/{slug}/v1/chat/completions",
                          "health: /api/vlm_serve/health — upstream /v1/models probe",
                          "클라이언트는 route_slug로 호출 → 모델 교체에 무관",
                      ])
    _add_footer(slide, "대표 파이프라인: UI-Venus(coarse) → MAI-UI(fine) → PaddleOCR-VL(검증)")


def slide_efficiency(prs):
    slide = _blank(prs)
    _add_title_band(slide, "효율 — 단일 거대 모델(Kimi-K2) 대비 약 20배",
                    "GUI grounding + OCR 작업 표면에 특화해 하드웨어를 크게 절감")
    cards = [
        ("필요 하드웨어", "H200 2장", "vs Kimi-K2 약 8장 → 약 4배 절감", NAVY),
        ("GPU당 모델 밀도", "2.5 모델 / H200", "vs 0.125 → 약 20배 우위", ACCENT),
        ("가중치 풋프린트", "~51 GiB", "vs ~1,040 GiB → 약 20배 절감", NAVY),
    ]
    top = Inches(1.6)
    height = Inches(2.3)
    width = Inches(4.0)
    gap = Inches(0.27)
    left = Inches(0.4)
    for role, big, sub, color in cards:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = color
        box.line.width = Pt(1.5)
        box.shadow.inherit = False
        rb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), width - Inches(0.5), Inches(0.4))
        run = rb.text_frame.paragraphs[0].add_run(); run.text = role
        _set_run(run, size=12.5, bold=True, color=color)
        bb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.75), width - Inches(0.5), Inches(0.8))
        run = bb.text_frame.paragraphs[0].add_run(); run.text = big
        _set_run(run, size=28, bold=True, color=NAVY)
        sb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(1.6), width - Inches(0.5), Inches(0.6))
        run = sb.text_frame.paragraphs[0].add_run(); run.text = sub
        _set_run(run, size=11.5, color=MUTED)
        left += width + gap
    _add_section_card(slide, left=Inches(0.4), top=Inches(4.15), width=Inches(12.55), height=Inches(2.7),
                      heading="해석", accent=ACCENT, bullets=[
                          "본 스택은 280 GiB 중 ~205 GiB로 5개 특화 모델 서빙 — 동일 2장으로 Kimi-K2는 1개도 적재 불가",
                          "레이턴시도 불리하지 않음: 7–8B dense VLM ~80–150 tok/s/req vs MoE 32B-active ~30–60",
                          "포기하는 것은 범용 추론 능력 — '실제로 읽어야 하는 화면'에 불필요한 범용성을 의도적으로 버린 트레이드오프",
                      ])
    _add_footer(slide, "근거: docs/setup_vlms/05-resource-comparison-vs-kimi-k2.md")


def slide_workflow1(prs):
    slide = _blank(prs)
    _add_title_band(slide, "workflow_1 — RCS GUI 자동화 + CCTV 캡처 PoC",
                    "VLM 기반 GUI 자동화 가능성 증명 (검증 후 동결)")
    steps = [
        ("Poll", "알람 API 1~2분 폴링, ALID=9006(Align Fail)만 필터"),
        ("Locate", "UI-Venus coarse bbox → MAI-UI 정밀 클릭점(2-stage)"),
        ("Verify", "PaddleOCR-VL로 입력값 OCR 재확인(closed-loop)"),
        ("Open DVR", "Tool DVR(CCTV) 자동 오픈, Channel 4 확대"),
        ("Capture", "최대 8분(~4,800 프레임) 100ms 간격 저장"),
    ]
    top = Inches(1.5)
    width = Inches(2.45)
    height = Inches(1.55)
    gap = Inches(0.07)
    left = Inches(0.4)
    for idx, (title, detail) in enumerate(steps, start=1):
        _add_pipeline_step(slide, left=left, top=top, width=width, height=height, index=idx, title=title, detail=detail)
        left += width + gap
    _add_section_card(slide, left=Inches(0.4), top=Inches(3.3), width=Inches(6.25), height=Inches(3.5),
                      heading="증명한 것", bullets=[
                          "VLM coarse→fine 좌표 인식이 UIA 없이 RCS UI를 안정적으로 클릭할 만큼 견고",
                          "Align Fail 시점 챔버 영상 100% 자동 보존 → 무인 모니터링",
                          "수집 프레임을 후속 VLM 자동 진단의 학습/검증 데이터로 재활용",
                          "DPI 보정·Tool 이름 OCR 정규화·edge-trigger 중복 제거",
                      ])
    _add_section_card(slide, left=Inches(6.85), top=Inches(3.3), width=Inches(6.1), height=Inches(3.5),
                      heading="한계 / 배운 점 → 동결", accent=ACCENT, bullets=[
                          "GUI 화면 읽기 ≠ align key 찾기 — 정밀 매칭은 CV 영역(→ wf2/wf3)",
                          "전 단계 VLM 의존은 비용·지연·edge-case 부담 → full 자동화 경제성 한계",
                          "production 경로는 workflow_3로 이전, CCTV/초기 실험만 잔류",
                          "align_images 데이터 루트 계약(rcp/msr/captured)의 원조",
                      ])
    _add_footer(slide, "Trigger: ALID=9006  •  2-stage VLM locator  •  production은 wf3로 이전")


def slide_workflow2(prs):
    slide = _blank(prs)
    _add_title_band(slide, "workflow_2 — 오프라인 CV 평가 벤치",
                    "golden set으로 matching/ensemble/consensus를 객관 검증·A/B·튜닝")
    _add_section_card(slide, left=Inches(0.4), top=Inches(1.4), width=Inches(6.25), height=Inches(2.6),
                      heading="3대 Golden Driver", bullets=[
                          "golden_localization — rcp 단독 localization",
                          "golden_consensus — consensus vs rcp A/B",
                          "golden_combined — production 라우팅 end-to-end + 3축 + OM/SEM 층화",
                          "지표: in_topk(recall), rank1(top 정답) · GT는 cond.txt + LOO",
                      ])
    _add_section_card(slide, left=Inches(6.85), top=Inches(1.4), width=Inches(6.1), height=Inches(2.6),
                      heading="핵심 CV 기법", accent=ACCENT, bullets=[
                          "Ensemble 3채널(Canny/Scharr/orientation+Chamfer) → RRF 융합 → NCC rerank",
                          "Youden 임계: match=0.6053, adjust=0.4727",
                          "Consensus: 최근 성공 S crop을 crosshair 기준 co-register 후 median blend",
                          "ensemble_lab로 production 무수정 실험(edge_ncc 등)",
                      ])
    _add_section_card(slide, left=Inches(0.4), top=Inches(4.15), width=Inches(12.55), height=Inches(2.7),
                      heading="역할·원칙", bullets=[
                          "production 엔진 무수정 — poc.workflow_3.align을 import만(반대 금지)",
                          "'가정 말고 측정' — 정책 변경 전 golden set 수치 우선 확인",
                          "벤치→프로덕션 파이프라인: 검증된 변경만 bit-parity 포팅 → 회귀 위험 통제",
                          "bench/config 분리 + [DIGEST] 한 줄 회신으로 오피스↔개발 전달 비용 최소화",
                      ])
    _add_footer(slide, "동결 아님 — 활성 연구·검증 harness")


def slide_consensus_lift(prs):
    slide = _blank(prs)
    _add_title_band(slide, "핵심 발견 — Consensus 재등록으로 정렬 정확도 도약",
                    "등록 align key(rcp) 단독의 천장을 뚫은 alignment 정확도 개선")
    metrics = [
        ("in_topk (proposer recall)", "0.434", "0.876", "+0.442 (약 +102%)"),
        ("rank1 (top 후보 정답)", "0.318", "0.764", "+0.446"),
    ]
    top = Inches(1.6)
    height = Inches(1.75)
    width = Inches(12.55)
    left = Inches(0.4)
    for name, before, after, delta in metrics:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        box.fill.solid(); box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = NAVY; box.line.width = Pt(1.0); box.shadow.inherit = False
        nb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.55), Inches(4.2), Inches(0.7))
        run = nb.text_frame.paragraphs[0].add_run(); run.text = name
        _set_run(run, size=15, bold=True, color=NAVY)
        b1 = slide.shapes.add_textbox(left + Inches(4.7), top + Inches(0.45), Inches(2.0), Inches(0.9))
        run = b1.text_frame.paragraphs[0].add_run(); run.text = before
        _set_run(run, size=30, bold=True, color=MUTED)
        _add_arrow(slide, left=left + Inches(6.75), top=top + Inches(0.65), width=Inches(0.7), height=Inches(0.45), color=ACCENT)
        b2 = slide.shapes.add_textbox(left + Inches(7.6), top + Inches(0.45), Inches(2.0), Inches(0.9))
        run = b2.text_frame.paragraphs[0].add_run(); run.text = after
        _set_run(run, size=30, bold=True, color=GREEN)
        b3 = slide.shapes.add_textbox(left + Inches(9.7), top + Inches(0.55), Inches(2.6), Inches(0.7))
        run = b3.text_frame.paragraphs[0].add_run(); run.text = delta
        _set_run(run, size=15, bold=True, color=ACCENT)
        top += height + Inches(0.25)
    _add_section_card(slide, left=Inches(0.4), top=Inches(5.5), width=Inches(12.55), height=Inches(1.35),
                      heading="해석 & 단서", accent=ACCENT, bullets=[
                          "같은 recipe의 최근 성공 외관을 따라가는 consensus가 stale rcp를 크게 능가 (golden set 벤치, LOO A/B, min_s=3)",
                          "주의: 벤치 기준 수치 — 오피스 실데이터 라우팅 종합·OM/SEM 층화는 golden_combined [DIGEST]로 확정 예정",
                      ])
    _add_footer(slide, "근거: poc/workflow_3/README.md · workflow_2 bench")


def slide_workflow3_loop(prs):
    slide = _blank(prs)
    _add_title_band(slide, "workflow_3 — 실시간 Align Fail 모니터링 루프 (현재 주력)",
                    "감지 → 접속 → CV 보정 → 실패 시 알림 → 상시 녹화 → 자동 종료")
    nodes = [
        ("1. 알람 감지", "ALID=9006 폴링\n(edge-trigger 중복 제거)", NAVY),
        ("2. RCS 접속", "tool 선택·접속\n+ 상시 녹화 시작", ACCENT),
        ("3. CV 보정", "consensus/rcp 라우팅\nmatch → reposition+OK", NAVY),
        ("4. 알림·녹화", "실패 시 cube 알림\n엔지니어 조작까지 녹화", ACCENT),
        ("5. 자동 종료", "engineer-done 감지 시\ntool 닫고 다음 대기", NAVY),
    ]
    top = Inches(1.7)
    height = Inches(1.7)
    width = Inches(2.3)
    gap = Inches(0.22)
    left = Inches(0.4)
    for idx, (title, detail, color) in enumerate(nodes):
        _add_loop_node(slide, left=left, top=top, width=width, height=height, title=title, detail=detail, color=color)
        if idx < len(nodes) - 1:
            _add_arrow(slide, left=left + width + Inches(0.01), top=top + (height - Inches(0.4)) / 2,
                       width=Inches(0.2), height=Inches(0.4), color=NAVY)
        left += width + gap
    fb = slide.shapes.add_textbox(Inches(0.4), top + height + Inches(0.12), Inches(12.55), Inches(0.35))
    p = fb.text_frame.paragraphs[0]; p.alignment = 2
    run = p.add_run(); run.text = "↻  다음 장비 대기 — 알람마다 반복 (cleanup은 try/finally로 보장)"
    _set_run(run, size=11.5, bold=True, color=ACCENT)
    _add_section_card(slide, left=Inches(0.4), top=Inches(4.4), width=Inches(6.25), height=Inches(2.45),
                      heading="회귀 위험 0 설계", bullets=[
                          "office 모듈 부재 시 자동 비활성 → 기존 동작 불변",
                          "consensus 부적격(부족/blur/cold)이면 modality별 rcp 폴백",
                          "ALIGN_FAIL_CONSENSUS=0 킬스위치 → 순수 rcp",
                          "실보정 2단계 게이트(SAFE_MODE=0 + DRY_RUN=0)",
                      ])
    _add_section_card(slide, left=Inches(6.85), top=Inches(4.4), width=Inches(6.1), height=Inches(2.45),
                      heading="부가 능력", accent=ACCENT, bullets=[
                          "상시 녹화: 변화 감지 적응 캡처(수동 조작 보존)",
                          "engineer-done: Recipe Monitor 카운터 hybrid 판독",
                          "feasibility: 모호 시 align key 재등록 권고",
                          "zoom ladder / PM dropdown: 배율 바꿔 재매칭",
                      ])
    _add_footer(slide, "monitor → {rcs, align, sem_monitor, runner, vlm, util}  •  4-layer DAG")


def slide_status(prs):
    slide = _blank(prs)
    _add_title_band(slide, "현재 상태 — 완료 vs 대기",
                    "코드는 완료, 오피스 활성화·캘리브레이션이 남은 단계")
    _add_section_card(slide, left=Inches(0.4), top=Inches(1.45), width=Inches(6.25), height=Inches(3.6),
                      heading="✅ 완료", accent=GREEN, bullets=[
                          "VLM 배포·운영 (H200×2, 5 모델)",
                          "workflow_1 GUI 자동화 + CCTV PoC (동결)",
                          "workflow_2 CV 평가 벤치 (활성)",
                          "wf3 루프·primary 보정·fallback search (코드)",
                          "consensus 라우팅·box-crop 템플릿 (코드)",
                          "재등록 플래깅·check-only 진단",
                      ])
    _add_section_card(slide, left=Inches(6.85), top=Inches(1.45), width=Inches(6.1), height=Inches(3.6),
                      heading="🟡 대기 / 진행 중", accent=ACCENT, bullets=[
                          "consensus 라이브 보정 활성화 (office_success_downloader 구현)",
                          "오피스 캘리브레이션 (zoom/click 좌표, engineer-done)",
                          "SEM-box landmark crop (모델별)",
                          "align_images 루트 이전 (wf1 → wf3)",
                          "golden_combined 오피스 실행 → 정확도 확정",
                          "pilot 실보정 (dry-run 후 단일 장비)",
                      ])
    _add_section_card(slide, left=Inches(0.4), top=Inches(5.2), width=Inches(12.55), height=Inches(1.65),
                      heading="리스크 완화", bullets=[
                          "downloader 부재→자동 비활성, cold cache→bounded sync 후 rcp 폴백, 실보정→2단계 게이트+pilot",
                          "정확도 수치는 벤치 기준 — 오피스 실데이터 [DIGEST]로 확정 예정",
                      ])
    _add_footer(slide, "상세: docs/project_progress/05_status_roadmap.md")


def slide_scalability(prs):
    slide = _blank(prs)
    _add_title_band(slide, "확장성 & 다음 단계",
                    "모델·장비·데이터·정확도 네 축의 확장 경로")
    _add_section_card(slide, left=Inches(0.4), top=Inches(1.45), width=Inches(6.25), height=Inches(3.6),
                      heading="확장성 (Scalability)", bullets=[
                          "모델: Flask proxy에 .env+모듈 1개로 6번째 추가(GPU 1 ~50 GiB 헤드룸), multi-GPU parallel",
                          "장비·recipe: consensus 이력 풀이 <class>/<recipe> 키(장비 무관) → 학습 데이터 공유",
                          "데이터 자산화: 상시 녹화 + recording_filter → interaction timeline → 모방학습",
                          "VLM ROI prior: align key 영역 grounding으로 CV 탐색 범위 축소(2-tier fallback)",
                      ])
    _add_section_card(slide, left=Inches(6.85), top=Inches(1.45), width=Inches(6.1), height=Inches(3.6),
                      heading="다음 주 우선순위", accent=ACCENT, bullets=[
                          "1. office_success_downloader 구현 → consensus 라이브 보정 활성화",
                          "2. 오피스 캘리브레이션(zoom/click, engineer-done) 완료",
                          "3. golden_combined 오피스 실행 → OM/SEM 층화·라우팅 정확도 확정",
                          "4. pilot 장비 1대 dry-run → 실보정 단계적 전환",
                      ])
    _add_section_card(slide, left=Inches(0.4), top=Inches(5.2), width=Inches(12.55), height=Inches(1.65),
                      heading="요약", accent=NAVY, bullets=[
                          "VLM 인프라 → GUI 자동화 → CV 정확도 → 실시간 통합 루프까지 PoC 완주",
                          "남은 것은 오피스 활성화·캘리브레이션과 pilot 검증 — 회귀 위험을 통제한 단계적 전개",
                      ])
    _add_footer(slide, "목적 · PoC 방향 · 성과 · 확장성  —  end")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_problem_direction(prs)
    slide_timeline(prs)
    slide_vlm_roles(prs)
    slide_efficiency(prs)
    slide_workflow1(prs)
    slide_workflow2(prs)
    slide_consensus_lift(prs)
    slide_workflow3_loop(prs)
    slide_status(prs)
    slide_scalability(prs)

    prs.save(OUTPUT_PATH)
    print(f"[INFO] PPTX saved: {OUTPUT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
