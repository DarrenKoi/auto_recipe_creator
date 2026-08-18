# -*- coding: utf-8 -*-
"""Align Tuning Agent 경과 보고 - PowerPoint(.pptx, 흰 배경) 생성 스크립트.

짝 스크립트 build_slides_bento.py 와 **같은 내용/같은 좌표계**를 쓴다. 16:9 슬라이드가
13.333 x 7.5 inch 이므로 96dpi 기준 1280x720 px 좌표가 1:1 로 대응한다(px * 9525 = EMU).
덕분에 두 판의 레이아웃을 따로 관리하지 않는다.

본편 4장 + 백업 2장이며, 백업은 show="0" 으로 숨겨 슬라이드쇼와 PDF 내보내기에서 빠진다.
차트는 이미지가 아니라 네이티브 PowerPoint 차트라 수치를 직접 편집할 수 있다.

한글 폰트는 오피스 Windows 기준 맑은 고딕(latin/ea/cs 모두 지정). Mac 에는 그 폰트가 없어
LibreOffice 렌더가 깨져 보이므로, 레이아웃 확인은 두 가지로 한다.
  1) 폰트 메트릭 실측 오버플로 검사 - 매 실행 시 자동(PIL + 시스템 한글 폰트가 있을 때).
  2) PREVIEW=1 - 폰트를 Apple SD Gothic Neo 로 바꾼 사본을 임시 폴더에 만든다. 이 사본을
     PDF/PNG 로 렌더해 눈으로 본다. 커밋 대상은 언제나 맑은 고딕 판이다.

실행:
  uv run python docs/project_progress/build_slides_pptx.py
  PREVIEW=1 uv run python docs/project_progress/build_slides_pptx.py   # Mac 확인용 사본
출력:
  docs/project_progress/Align_Tuning_Agent.pptx
"""
import re
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

import os
import tempfile
from pathlib import Path

HERE = Path(__file__).resolve().parent
# PREVIEW=1 은 Mac 렌더 확인 전용 사본이다. 커밋 대상(맑은 고딕 판)을 덮어쓰지 않는다.
PREVIEW = os.environ.get("PREVIEW") == "1"
OUT = (Path(tempfile.gettempdir()) / "Align_Tuning_Agent_preview.pptx") if PREVIEW \
    else (HERE / "Align_Tuning_Agent.pptx")
FONT = "Apple SD Gothic Neo" if PREVIEW else "맑은 고딕"
# 오버플로 검사용 프록시 폰트(맑은 고딕과 폭이 비슷한 한글 sans). 없으면 검사를 건너뛴다.
METRIC_FONT = "/System/Library/Fonts/AppleSDGothicNeo.ttc"
try:
    from PIL import ImageFont
    METRIC_AVAILABLE = os.path.exists(METRIC_FONT)
except ImportError:
    METRIC_AVAILABLE = False

C = lambda h: RGBColor.from_string(h)
BG      = "FFFFFF"
TEXTC   = "10233A"
MUTED   = "56677D"
DIM     = "8A97A8"
ACC     = "0E6FB8"
ACC_DK  = "0A5389"
SOFT    = "EAF3FB"   # 강조 카드 배경
CARD    = "F6F9FC"   # 일반 카드 배경
BORDER  = "DCE5EE"
BORDER2 = "B9D6EC"
NEUT    = "AEBECD"   # As-Is 막대

_BOXES = []
_SI = [0]
PXE = 9525
px = lambda v: Emu(int(round(v * PXE)))
tp = lambda v: Pt(round(v * 0.75, 1))


def _font(run, size, bold, color, name=FONT):
    f = run.font
    f.size = tp(size)
    f.bold = bold
    f.color.rgb = C(color)
    f.name = name
    rPr = f._element
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def T(sl, x, y, w, h, html, size=16, weight=600, color=TEXTC, align="left",
      valign="top", lh=1.45):
    _BOXES.append((_SI[0], x, y, w, h, html, size, lh))
    tb = sl.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[valign]
    html = html.replace("&nbsp;", " ")
    for i, line in enumerate(re.split(r"<br\s*/?>", html)):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        p.line_spacing = lh
        p.space_after = Pt(0)
        for chunk in re.split(r"(<b>.*?</b>)", line):
            if not chunk:
                continue
            b = chunk.startswith("<b>")
            txt = re.sub(r"</?b>", "", chunk)
            r = p.add_run()
            r.text = txt
            _font(r, size, True if b else (weight >= 700), color)
    return tb


def R(sl, x, y, w, h, fill, border=None, bw=1, radius=12):
    shp = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        px(x), px(y), px(w), px(h))
    if radius:
        shp.adjustments[0] = min(0.5, radius / float(min(w, h)))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = C(fill)
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = C(border)
        shp.line.width = Pt(bw * 0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def LINE(sl, x, y, w, color=ACC, bw=1.5, dash=True):
    cn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x), px(y), px(x + w), px(y))
    cn.line.color.rgb = C(color)
    cn.line.width = Pt(bw)
    if dash:
        cn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return cn


def VLINE(sl, x, y, h, color=BORDER, bw=1.25):
    cn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x), px(y), px(x), px(y + h))
    cn.line.color.rgb = C(color)
    cn.line.width = Pt(bw)
    return cn


def TABLE(sl, x, y, w, h, cols, rows, fs=13.5, padY=8, hdr_fill=SOFT, hdr_color=ACC):
    g = sl.shapes.add_table(len(rows), len(cols), px(x), px(y), px(w), px(h))
    t = g.table
    t.first_row = True
    t.horz_banding = False
    tot = float(sum(cols))
    for i, cw in enumerate(cols):
        t.columns[i].width = px(w * cw / tot)
    rh = h / float(len(rows))
    for r_i, row in enumerate(rows):
        t.rows[r_i].height = px(rh)
        for c_i, cell in enumerate(row):
            spec = cell if isinstance(cell, dict) else {"html": cell}
            cl = t.cell(r_i, c_i)
            cl.margin_left = cl.margin_right = px(12)
            cl.margin_top = cl.margin_bottom = px(padY)
            cl.vertical_anchor = MSO_ANCHOR.MIDDLE
            cl.fill.solid()
            cl.fill.fore_color.rgb = C(hdr_fill if r_i == 0 else
                                       (BG if r_i % 2 else "FAFCFE"))
            tf = cl.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.line_spacing = 1.3
            col = hdr_color if r_i == 0 else spec.get("color", TEXTC)
            bold = r_i == 0 or spec.get("bold", False)
            for chunk in re.split(r"(<b>.*?</b>)", spec["html"]):
                if not chunk:
                    continue
                run = p.add_run()
                run.text = re.sub(r"</?b>", "", chunk)
                _font(run, fs, bold or chunk.startswith("<b>"), col)
    return t


prs = Presentation()
prs.slide_width = px(1280)
prs.slide_height = px(720)
BLANK = prs.slide_layouts[6]


def new():
    _SI[0] += 1
    sl = prs.slides.add_slide(BLANK)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = C(BG)
    return sl


def notes(sl, txt):
    sl.notes_slide.notes_text_frame.text = txt


STEPS = [
    ("01", "VLM 인프라 확보", "사내 HCP에 오픈소스 소형 모델 배포\n GPU 1~2장 규모 · 통합 API 운영", "완료 · 2종 상시 운영"),
    ("02", "GUI 자동화 증명", "비표준 계측 화면에서 좌표 인식\n coarse → fine 2단계 + OCR 확인", "완료 · 정확도 90%+"),
    ("03", "정렬 정확도 확보", "오프라인 벤치에서 A/B·튜닝\n 검증 통과한 방식만 운영 반영", "완료 · 운영 반영"),
    ("04", "실시간 통합 루프", "감지 → 접속 → 보정 → 알림 → 녹화\n end-to-end 무인 루프", "진행 · 실운전 안정화"),
]

# ══════════════════════════════════════════════ S1 · Head Message
s = new()
R(s, 0, 0, 1280, 8, ACC, radius=0)
T(s, 96, 84, 700, 30, "과제 종류 · 현업&nbsp;&nbsp;|&nbsp;&nbsp;경과 보고", 19, 700, ACC)
T(s, 96, 124, 1000, 108, "Align Tuning Agent", 84, 900, TEXTC, lh=1.05)
R(s, 96, 244, 110, 6, ACC, radius=3)
T(s, 96, 276, 1010, 72,
  "Align Fail 무인 대응 AI Agent&nbsp;&nbsp;·&nbsp;&nbsp;화면을 이해하는 VLM과 좌표를 결정하는 CV의 역할 분리 설계",
  21, 600, MUTED, lh=1.5)
T(s, 96, 344, 1088, 84,
  "수동 <b>5분</b> 대응을 무인 <b>1분</b> 대응으로 — 회수한 장비 가동 시간으로 <b>주 WAFER 100장</b>을 더 측정합니다.",
  25, 800, TEXTC, lh=1.4)
for n, (v, lab) in enumerate([
        ("4분", "1건당 대응 지연 단축<br>감지~보정 5분 → 1분"),
        ("주 100장", "추가 측정 가능한 WAFER<br>400건 × 4분 ÷ 15분/장"),
        ("28대", "대상 계측 장비<br>R3 CD-SEM")]):
    x = 96 + n * 374
    R(s, x, 432, 340, 134, SOFT, BORDER2, 1, 14)
    T(s, x + 22, 450, 296, 58, v, 40, 900, ACC)
    T(s, x + 22, 514, 296, 46, lab, 13.5, 600, MUTED, lh=1.5)
R(s, 96, 594, 1088, 78, ACC, radius=12)
T(s, 120, 610, 400, 34, "예상 종료 · <b>2026년 9월 말</b>", 23, 800, "FFFFFF")
T(s, 520, 612, 640, 50,
  "현업 엔지니어 합동 실전 테스트와 단일 장비 Pilot 완료 시점<br>이후 4Q 대상 장비 확대 · 2027년 Recipe Tuning 확장",
  13.5, 600, "D6E9F7", lh=1.5)
notes(s, "표지 겸 Head Message. 주 400건의 Align Fail에 대해 1건당 대응 지연을 5분에서 1분으로 줄여 "
         "주 WAFER 약 100장의 측정 여력을 회수한다. 대상은 R3 CD-SEM 28대. "
         "환산 근거: 400건/주 x 4분 단축 / 15분/장 = 주 약 100장. 예상 종료는 2026년 9월 말.")

# ══════════════════════════════════════════════ S2 · 선정 배경과 목적
s = new()
T(s, 96, 50, 900, 54, "선정 배경과 목적", 36, 900, TEXTC)
T(s, 96, 110, 1088, 30,
  "Align Fail이 나면 장비가 멈추고, Line 구성원이 직접 접속해 정렬을 다시 맞춰야 재개됩니다.", 16, 600, MUTED)
T(s, 96, 150, 520, 26, "AS-IS · 현재의 문제", 14, 800, ACC)
for n, (v, lab) in enumerate([
        ("주 400건", "R3 CD-SEM 28대에서 발생하는 Align Fail 알람"),
        ("1건당 5분", "Line 구성원의 감지 → RCS 접속 → 수동 재정렬"),
        ("주 33시간", "누적 대응 지연 — 그만큼 장비가 멈춰 있습니다")]):
    y = 186 + n * 86
    R(s, 96, y, 520, 74, CARD, BORDER, 1, 10)
    T(s, 116, y + 16, 180, 40, v, 28, 900, TEXTC)
    T(s, 300, y + 18, 300, 44, lab, 13, 600, MUTED, lh=1.45)
T(s, 96, 454, 520, 26, "왜 자동화가 어려웠는가", 14, 800, ACC)
for n, txt in enumerate([
        "<b>표준 컴포넌트가 아닌 계측 화면</b> — 일반 UI 자동화 도구가 버튼·입력창을 인식하지 못해, 화면을 이미지로 보고 좌표를 찾아야 합니다.",
        "<b>되풀이될 수밖에 없는 현상</b> — 등록 Align Key와 실제 Wafer 이미지가 공정 Variation으로 벌어집니다. Split Lot이 많은 R3에서는 상시 발생합니다.",
        "<b>의미 이해와 픽셀 정밀도를 동시에</b> — 매번 조금씩 달라지는 Align Key를 찾아 정확한 좌표로 재정렬해야 합니다."]):
    y = 486 + n * 62
    T(s, 96, y + 5, 14, 20, "•", 13, 900, ACC)
    T(s, 116, y, 500, 58, txt, 12.5, 600, MUTED, lh=1.5)
T(s, 656, 150, 528, 26, "TO-BE · 과제 목표", 14, 800, ACC)
TABLE(s, 656, 184, 528, 338, [0.72, 1.0, 1.35],
      [["구분", "As-Is (수동)", "To-Be (무인 Agent)"],
       ["감지", "사람이 수시로 확인", "10초 주기 자동 확인 · 9006만 선별"],
       ["접속·보정", "엔지니어가 RCS 접속 후 수동 재정렬", "RCS 자동 접속 → 성공 사례 기반 자동 보정"],
       ["실패·오측", "사람이 발견할 때까지 지연", "담당자 Cube 알림 · 오측 시 측정 중단"],
       ["기록", "별도 기록 없음", "상시 녹화 + 성공/실패 이미지 자동 저장"],
       [{"html": "1건당 시간", "bold": True}, {"html": "평균 5분", "bold": True},
        {"html": "목표 1분", "color": ACC, "bold": True}]], fs=12.5, padY=7)
R(s, 656, 540, 528, 132, SOFT, BORDER2, 1, 12)
T(s, 680, 558, 480, 26, "과제 목적", 13.5, 800, ACC)
T(s, 680, 588, 480, 72,
  "사람이 하던 <b>감지 → 장비 접속 → 정렬 재수정</b>을 AI가 대신 맡아 장비 idle 시간을 줄이고 계측 처리량을 회복합니다.",
  15, 600, TEXTC, lh=1.55)
notes(s, "주 400건 x 1건당 5분 = 주 33시간의 장비 정지. 어려웠던 이유 3가지 - 비표준 UI, "
         "공정 Variation으로 인한 이미지 괴리(R3는 Split Lot이 많아 상시), 의미 이해와 픽셀 정밀도의 동시 요구.")

# ══════════════════════════════════════════════ S3 · 추진 전략과 성과
s = new()
T(s, 96, 44, 900, 54, "추진 전략과 성과", 36, 900, TEXTC)
T(s, 96, 102, 1088, 26,
  "네 단계로 나눠 쌓았습니다. 앞 단계의 검증을 통과한 방식만 다음 단계에 반영했습니다.", 14, 600, MUTED)
LINE(s, 96, 140, 1088)
for n, (num, ttl, desc, chip) in enumerate(STEPS):
    x = 96 + n * 278
    R(s, x, 152, 254, 152, CARD, BORDER2 if n == 3 else BORDER, 1, 12)
    T(s, x + 18, 162, 60, 36, num, 24, 900, ACC)
    T(s, x + 18, 198, 218, 28, ttl, 16, 800, TEXTC)
    T(s, x + 18, 228, 218, 48, desc.replace("\n", "<br>"), 11.5, 600, MUTED, lh=1.5)
    T(s, x + 18, 278, 218, 22, chip, 11, 700, ACC if n < 3 else ACC_DK)
T(s, 96, 328, 560, 24, "정량적 성과 · 정렬 위치 재조정 정확도", 14, 800, ACC)
T(s, 96, 352, 560, 22,
  "Recipe 200개 오프라인 벤치 · 등록 이미지 1장 → 최근 성공 사례 종합(consensus)", 11.5, 600, DIM)
cd = CategoryChartData()
cd.categories = ["후보 포함률 (in_topk)", "1순위 적중률 (rank1)"]
cd.add_series("등록 이미지 1장", (0.434, 0.318))
cd.add_series("성공 사례 종합(consensus)", (0.876, 0.764))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, px(88), px(376), px(576), px(284), cd)
ch = gf.chart
ch.has_legend = True
ch.legend.position = XL_LEGEND_POSITION.TOP
ch.legend.include_in_layout = False
ch.font.size = tp(13)
ch.font.name = FONT
ch.font.color.rgb = C(MUTED)
pl = ch.plots[0]
pl.gap_width = 80
pl.overlap = -10
pl.has_data_labels = True
pl.data_labels.number_format = "0.000"
pl.data_labels.number_format_is_linked = False
pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
pl.data_labels.font.size = tp(12)
pl.data_labels.font.bold = True
pl.data_labels.font.color.rgb = C(TEXTC)
for i, col in enumerate((NEUT, ACC)):
    sr = pl.series[i]
    sr.format.fill.solid()
    sr.format.fill.fore_color.rgb = C(col)
    sr.format.line.fill.background()
va = ch.value_axis
va.minimum_scale = 0.0
va.maximum_scale = 1.0
va.has_major_gridlines = True
va.format.line.color.rgb = C(BORDER)
ch.category_axis.format.line.color.rgb = C(BORDER)
T(s, 96, 662, 576, 24,
  "후보 포함률 <b>0.434 → 0.876</b> (약 2배) · 1순위 적중률 <b>0.318 → 0.764</b> (2배 이상)",
  12.5, 700, TEXTC)
T(s, 688, 328, 496, 24, "정성적 성과 · 구조적으로 확보한 것", 14, 800, ACC)
for n, (ttl, desc) in enumerate([
        ("VLM 인프라 자체 확보", "GPU 1~2장 규모 소형 모델을 사내 HCP에서 직접 구동. 벤치로 이긴 2종만 상시 운영해 모델 교체를 측정으로 결정합니다."),
        ("실장비 보정 반자동 개통 (2026-08)", "정렬 위치 재조정은 자동, 마지막 확정(OK)은 엔지니어. 개통 과정에서 보정이 한 번도 실행되지 않던 병목을 찾아 제거했습니다."),
        ("암묵지 자산화 경로 구축", "엔지니어 조작 녹화 → 조작 타임라인 → 한국어 절차서. 입력은 후킹하지 않고 화면만 관측해 동의 범위를 코드에 고정했습니다.")]):
    y = 356 + n * 104
    R(s, 688, y, 496, 94, CARD, BORDER, 1, 10)
    T(s, 708, y + 12, 456, 24, ttl, 14, 800, TEXTC)
    T(s, 708, y + 38, 456, 46, desc, 11.5, 600, MUTED, lh=1.5)
T(s, 688, 664, 496, 22,
  "※ 모드별 정확도 · 재순위 채택 · 잔여 개선 분해는 백업 슬라이드 참조", 11, 700, DIM, align="right")
notes(s, "정량 핵심: 등록 이미지 1장 대비 최근 성공 사례를 종합하면 후보 포함률 0.434->0.876, "
         "1순위 적중률 0.318->0.764. 운영에서는 결국 한 곳으로 재조정하므로 1순위 적중률이 실제 성능에 "
         "가장 가깝다. 단, 오프라인 벤치 기준이며 오피스 실데이터 확정은 진행 중.")

# ══════════════════════════════════════════════ S4 · 기대 효과와 향후 일정
s = new()
T(s, 96, 44, 900, 54, "기대 효과와 향후 일정", 36, 900, TEXTC)
for n, (num, ttl, _d, _c) in enumerate(STEPS):
    x = 96 + n * 278
    R(s, x, 108, 254, 48, CARD, BORDER, 1, 8)
    T(s, x + 14, 120, 30, 24, num, 13, 900, ACC)
    T(s, x + 46, 120, 196, 24, ttl, 12.5, 700, MUTED)
T(s, 96, 172, 560, 26, "기대 효과", 14, 800, ACC)
TABLE(s, 96, 204, 560, 208, [1.35, 0.8],
      [["항목", "값"],
       ["대상 장비 수", "28대"],
       ["주 평균 Align Fail 발생", "약 400건"],
       ["1건당 대응 지연", {"html": "5분 → 1분", "color": ACC, "bold": True}],
       ["주간 절감 환산", {"html": "WAFER 약 100장 추가 측정", "color": ACC, "bold": True}]],
      fs=13, padY=8)
for n, txt in enumerate([
        "야간·주말 무인 감시·보정",
        "실패 시점 화면 증거 자동 보존 → 원인 분석 시간 단축",
        "수집 데이터의 VLM 학습 자산화",
        "같은 Recipe면 장비가 달라도 성공 이력 공유 → 확대 비용 최소"]):
    x = 96 + (n % 2) * 288
    y = 428 + (n // 2) * 74
    R(s, x, y, 272, 64, CARD, BORDER, 1, 8)
    T(s, x + 14, y + 12, 244, 44, txt, 11.5, 650, MUTED, lh=1.45)
T(s, 688, 172, 496, 26, "향후 일정", 14, 800, ACC)
for n, (when, what, hot) in enumerate([
        ("8월 후반", "첫 실알람 완주 확인 (엔지니어 입회) · 엔지니어 작업 완료 감지 판독 확인 · 녹화 1건의 절차 추출 판정", True),
        ("9월", "MI 현업 엔지니어와 합동 실전 테스트 · 단일 장비 Pilot · 엔지니어 PC 배포 방식 확정", True),
        ("4Q", "Pilot 결과를 반영해 대상 장비 확대", False),
        ("2027년", "VLM이 Recipe Tuning을 수행하도록 workflow 구조 확장", False)]):
    y = 204 + n * 74
    R(s, 692, y + 20, 12, 12, ACC if hot else NEUT, radius=6)
    if n < 3:
        VLINE(s, 698, y + 36, 54)
    T(s, 720, y + 12, 110, 26, when, 15, 800, ACC if hot else MUTED)
    T(s, 838, y + 12, 346, 56, what, 11.5, 600, MUTED, lh=1.45)
R(s, 688, 500, 496, 84, CARD, BORDER, 1, 10)
T(s, 708, 512, 456, 22, "가장 가까운 관문 — 개통과 완주는 다릅니다", 13, 800, TEXTC)
T(s, 708, 536, 456, 42,
  "실보정은 열렸지만 실알람 완주 사례가 아직 없습니다. 예상 허들과 진행 현황은 백업 슬라이드 참조.",
  11.5, 600, MUTED, lh=1.45)
R(s, 96, 606, 1088, 74, ACC, radius=12)
T(s, 120, 622, 400, 32, "예상 종료 · <b>2026년 9월 말</b>", 21, 800, "FFFFFF")
T(s, 520, 622, 640, 48,
  "현업 합동 실전 테스트와 단일 장비 Pilot 완료를 과제 종료 시점으로 봅니다.<br>이후 4Q 확대 적용은 후속 과제로 이어집니다.",
  12.5, 600, "D6E9F7", lh=1.5)
notes(s, "정량 효과 재확인(28대 / 주 400건 / 5분->1분 / 주 WAFER 100장)과 정성 효과 4가지. "
         "일정은 8월 후반 첫 실알람 완주 확인, 9월 합동 실전 테스트와 단일 장비 Pilot, 4Q 확대, "
         "2027년 Recipe Tuning 확장.")

# ══════════════════════════════════════ 백업 1 · 성과 상세 (숨김)
s = new()
T(s, 96, 54, 900, 48, "[백업] 성과 상세", 30, 900, TEXTC)
T(s, 96, 104, 1088, 26, "촬영 모드별 정확도와, 남은 개선 여지가 어디에 있는지에 대한 결론입니다.", 14, 600, MUTED)
TABLE(s, 96, 150, 560, 172, [1.0, 0.85, 0.85, 0.55],
      [["촬영 모드 (성공 사례 종합 기준)", "후보 포함률", "1순위 적중률", "표본"],
       ["OM · 저배율 (키가 화면의 10~20%)", "0.911", "0.852", "135"],
       ["SEM · 고배율 (키가 화면의 80~100%)", "0.789", "0.665", "185"]], fs=12.5, padY=9)
T(s, 96, 334, 560, 48,
  "통념과 달리 고배율 SEM이 더 까다롭습니다. 화면 대부분이 비슷한 반복 패턴이라 진짜 정렬 위치와 헷갈리는 가짜 후보가 많습니다.",
  12, 600, MUTED, lh=1.5)
for n, (v, lab) in enumerate([
        ("+0.073", "촬영 모드별 재순위 채택 시 SEM 1순위 적중률 상승 (OM +0.021, 무손실 확인)"),
        ("0.826", "재순위 반영 후 최종 채택 구성의 1순위 적중률"),
        ("1.000", "장비 창 버튼 판독 정확도 — 모델 벤치로 5종 평가 후 2종 상시 운영으로 정리")]):
    y = 150 + n * 92
    R(s, 688, y, 496, 80, SOFT, BORDER2, 1, 10)
    T(s, 708, y + 14, 150, 40, v, 26, 900, ACC)
    T(s, 862, y + 14, 302, 56, lab, 11.5, 600, MUTED, lh=1.45)
R(s, 96, 416, 1088, 120, SOFT, BORDER2, 1, 12)
T(s, 120, 434, 1040, 26, "잔여 개선 여지의 분해 — 알고리즘 축은 소진됐습니다", 16, 800, ACC)
T(s, 120, 464, 1040, 62,
  "재순위 반영 후 남은 실패의 <b>약 87%</b>가 '정답이 애초에 후보 목록에 없어 어떤 재순위로도 고칠 수 없는' 경우였습니다. "
  "남은 몫은 정렬 키를 더 잘 구별되는 위치로 <b>재등록</b>하는 일이며, 이는 현업 엔지니어와의 협업 과제입니다. "
  "이 결론에 따라 무게중심을 실시간 루프의 실운전 안정화로 옮겼습니다.", 13, 600, TEXTC, lh=1.6)
T(s, 96, 556, 1088, 50,
  "※ 위 정확도는 모두 오프라인 벤치(다운로드 샘플) 기준입니다. 오피스 실데이터 확정은 진행 중이며, 확정 전에는 성공 사례를 더 많이 요구하는 보수적 기준으로 운영합니다.",
  12, 600, DIM, lh=1.5)
notes(s, "백업. SEM이 구조적으로 더 어렵고, 남은 개선의 87%는 매칭 알고리즘이 아니라 정렬 키 재등록으로 풀어야 한다.")

# ══════════════════════════════════════ 백업 2 · 진행 현황과 예상 허들 (숨김)
s = new()
T(s, 96, 50, 900, 48, "[백업] 진행 현황과 예상 허들", 30, 900, TEXTC)
TABLE(s, 96, 108, 560, 452, [1.0, 1.0],
      [["항목", "상태"],
       ["VLM 배포·운영 (사내 HCP)", "완료 · 2종 상시 운영"],
       ["GUI 자동화 · 화면 증거 캡처", "완료 · 오피스 검증"],
       ["CV 정렬 정확도 (consensus + 모드별 재순위)", "완료 · 운영 반영"],
       ["실시간 루프 골격·보정·녹화·알림", "완료"],
       ["실장비 보정 (반자동, 확정은 사람)", {"html": "개통 · 완주 대기", "color": ACC, "bold": True}],
       ["루프 실패 경로 하드닝", "완료"],
       ["성공 사례 기반 보정 활성화", {"html": "이미지 공급 대기", "color": ACC, "bold": True}],
       ["지식 자산화 (녹화 → 절차서)", {"html": "녹화 확보 · 추출 미실행", "color": ACC, "bold": True}],
       ["엔지니어 PC 배포 방식", {"html": "방식 결정 대기", "color": ACC, "bold": True}],
       ["현업 합동 실전 테스트", {"html": "9월 예정", "color": ACC, "bold": True}]], fs=12, padY=6)
T(s, 688, 108, 496, 26, "예상 허들과 대응", 14, 800, ACC)
for n, (ttl, desc) in enumerate([
        ("개통과 완주는 다릅니다", "실알람 완주 사례가 없습니다. → replay 절차와 점검 모드로 단계별 사전 확인, 첫 실알람은 엔지니어 입회하에 관찰합니다."),
        ("오프라인 수치의 현장 검증 갭", "정확도는 벤치 기준입니다. → 오피스 실데이터로 확정하며, 그전까지는 보수적 기준으로 운영합니다."),
        ("성공 이미지 적재가 선행 조건", "적재 전에는 등록 이미지 폴백으로만 동작합니다. → 수집·적재 기능을 최우선 활성화 조건으로 진행합니다."),
        ("재등록은 협업 과제", "구별력이 약한 정렬 키는 알고리즘으로 풀리지 않습니다. → 우선순위 리스트로 순위화해 현업과 재등록을 협의합니다."),
        ("VLM 서버 가용성·배포 마찰", "재시작 시 환경 재설치에 3~4시간. → 실시간 루프는 호출이 적은 경로로 설계하고, 배포는 더블클릭 실행 방식을 사내 자원으로 고정합니다.")]):
    y = 140 + n * 86
    R(s, 688, y, 496, 76, CARD, BORDER, 1, 9)
    T(s, 706, y + 10, 460, 22, ttl, 13, 800, TEXTC)
    T(s, 706, y + 32, 460, 40, desc, 11, 600, MUTED, lh=1.4)
T(s, 96, 576, 560, 60,
  "실보정 전환은 <b>관찰 전용 → 반자동(현재) → 단일 장비 Pilot → 확대</b> 순서의 단계적 활성화로 리스크를 관리합니다.",
  12.5, 600, MUTED, lh=1.55)
notes(s, "백업. 진행 현황은 구현/오피스 검증/실알람 경험을 구분해 표기했다. 가장 가까운 관문은 첫 실알람 완주 확인.")

# 백업 2장은 슬라이드쇼에서 숨김 (본편 4장 유지)
for sl in list(prs.slides)[4:]:
    sl._element.set("show", "0")

# ── 폰트 메트릭 실측으로 텍스트 박스 오버플로 검사 (Apple SD Gothic Neo 프록시 +5%)
_FCACHE = {}
def _fnt(sz):
    k = int(round(sz))
    if k not in _FCACHE:
        _FCACHE[k] = ImageFont.truetype(METRIC_FONT, k)
    return _FCACHE[k]
def _lines(txt, size, boxw):
    f = _fnt(size); n = 1; cur = 0.0
    for chpos, ch in enumerate(txt):
        cw = f.getlength(ch) * 1.05
        if cur + cw > boxw and cur > 0:
            n += 1; cur = cw
        else:
            cur += cw
    return n
over = []
for si, x, y, w, h, html, size, lh in (_BOXES if METRIC_AVAILABLE else []):
    total = 0
    for seg in re.split(r"<br\s*/?>", html):
        t = re.sub(r"</?b>", "", seg).replace("&nbsp;", " ")
        total += _lines(t, size, w)
    need = total * size * lh
    if need > h + 1.5:
        over.append((si, y, round(need - h, 1), total, (re.sub(r"<[^>]+>", "", html)[:46])))
if not METRIC_AVAILABLE:
    print("[WARNING] PIL 또는 한글 프록시 폰트가 없어 오버플로 검사를 건너뜁니다:", METRIC_FONT)
else:
    print("[INFO] 텍스트 오버플로 검사 (slide, y, 초과px, 줄수, 앞부분)")
    for o in sorted(over):
        print("[WARNING]  ", o)
    if not over:
        print("[INFO]   없음")

prs.save(OUT)
print("[INFO] saved:", OUT)
print("[INFO] slides:", len(prs.slides._sldIdLst), "| 본편 4 + 백업 2(숨김)")
