"""Boss 보고용 2-slide PPTX 생성 스크립트.

poc/workflow_1/ 의 Align Fail 모니터링 + CCTV 자동 캡처 작업을 요약한다.
실행:
  uv run python poc/workflow_1/build_report_pptx.py
출력:
  poc/workflow_1/align_fail_capture_report.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUTPUT_PATH = Path(__file__).resolve().parent / "align_fail_capture_report.pptx"

NAVY = RGBColor(0x14, 0x2A, 0x55)
ACCENT = RGBColor(0xE8, 0x6A, 0x1F)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
TEXT = RGBColor(0x22, 0x29, 0x33)
MUTED = RGBColor(0x55, 0x60, 0x70)


def _set_run(run, *, size, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _add_title_band(slide, title_text, subtitle_text):
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.55)
    )
    tf = title_box.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    _set_run(run, size=26, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    sub_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.7), Inches(12.5), Inches(0.35)
    )
    p = sub_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = subtitle_text
    _set_run(run, size=13, color=RGBColor(0xCD, 0xD7, 0xE6))


def _add_section_card(slide, *, left, top, width, height, heading, bullets):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = NAVY
    card.line.width = Pt(0.75)
    card.shadow.inherit = False

    head_box = slide.shapes.add_textbox(
        left + Inches(0.25), top + Inches(0.15), width - Inches(0.5), Inches(0.45)
    )
    p = head_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = heading
    _set_run(run, size=16, bold=True, color=NAVY)

    body_box = slide.shapes.add_textbox(
        left + Inches(0.25),
        top + Inches(0.65),
        width - Inches(0.5),
        height - Inches(0.8),
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = f"• {bullet}"
        _set_run(run, size=12, color=TEXT)


def _add_pipeline_step(slide, *, left, top, width, height, index, title, detail):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = ACCENT
    box.line.width = Pt(1.25)
    box.shadow.inherit = False

    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + Inches(0.15),
        top + Inches(0.15),
        Inches(0.4),
        Inches(0.4),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    btf = badge.text_frame
    btf.margin_left = btf.margin_right = btf.margin_top = btf.margin_bottom = 0
    p = btf.paragraphs[0]
    p.alignment = 2  # center
    run = p.add_run()
    run.text = str(index)
    _set_run(run, size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

    title_box = slide.shapes.add_textbox(
        left + Inches(0.65), top + Inches(0.12), width - Inches(0.8), Inches(0.4)
    )
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, size=13, bold=True, color=NAVY)

    body_box = slide.shapes.add_textbox(
        left + Inches(0.2),
        top + Inches(0.6),
        width - Inches(0.4),
        height - Inches(0.7),
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = detail
    _set_run(run, size=10.5, color=MUTED)


def _add_footer(slide, text):
    box = slide.shapes.add_textbox(
        Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.3)
    )
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run(run, size=10, color=MUTED)


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_title_band(
        slide,
        "CD-SEM Align Fail 자동 감지 & 영상 캡처 파이프라인",
        "poc/workflow_1 — Align Fail 발생 시점의 CCTV 화면을 자동으로 보존",
    )

    _add_section_card(
        slide,
        left=Inches(0.4),
        top=Inches(1.4),
        width=Inches(6.25),
        height=Inches(2.7),
        heading="문제 정의",
        bullets=[
            "CD-SEM 장비에서 Align Fail (ALID=9006) 발생 시,"
            " 원인 분석을 위한 CCTV 영상이 사람이 늦게 확인하면 덮어쓰기로 사라짐",
            "엔지니어가 알람 콘솔을 수동 모니터링하고,"
            " EQP_ID 별로 Tool DVR 을 직접 열어야 했음",
            "Align Fail 순간의 챔버 내부 상태(Channel 4 영상)를"
            " 사후에 재현하기 어려움",
        ],
    )

    _add_section_card(
        slide,
        left=Inches(6.85),
        top=Inches(1.4),
        width=Inches(6.1),
        height=Inches(2.7),
        heading="이번에 만든 것",
        bullets=[
            "Align Fail 알람을 1~2분 주기로 자동 폴링 (ALID=9006 필터)",
            "감지 즉시 Windows 팝업 + 누적 텍스트 로그로 엔지니어에게 통보",
            "동일 EQP_ID 의 중복 알람은 edge-trigger 처리 (해제 후 재발 시에만 재알림)",
            "감지된 EQP_ID 의 Tool DVR(CCTV) 창을 자동으로 열고,"
            " Channel 4 를 확대한 뒤 최대 8분간 100ms 간격으로 프레임 저장",
        ],
    )

    _add_section_card(
        slide,
        left=Inches(0.4),
        top=Inches(4.25),
        width=Inches(12.55),
        height=Inches(2.7),
        heading="기대 효과",
        bullets=[
            "Align Fail 발생 시점의 챔버 영상을 100% 자동 보존 →"
            " 원인 분석/리포트 작성 시간 단축",
            "야간/주말 무인 모니터링 가능 — 엔지니어가 알람 콘솔을 지킬 필요 없음",
            "수집된 8분 분량 (~4,800 프레임) 데이터를"
            " 후속 VLM(UI-Venus) 기반 자동 진단의 학습/검증 데이터로 활용",
            "기존 RCS/DVR 시스템에 변경 없이 GUI 자동화로 동작 — 운영 리스크 최소화",
        ],
    )

    _add_footer(slide, "Module: poc/workflow_1/  •  Trigger: ALID=9006 (Align Fail)  •  Capture: CH4 @ 100ms, up to 8 min")


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_title_band(
        slide,
        "End-to-End 파이프라인 & 구성 모듈",
        "감지 → 알림 → CCTV 진입 → CH4 확대 → 프레임 캡처",
    )

    steps = [
        (
            "Poll",
            "CD-SEM 알람 API 를 1~2분 주기로 조회하여 ALID=9006 (Align Fail) 만 필터",
        ),
        (
            "Notify",
            "Windows MessageBox 팝업 + logs/align_fail_alarms.txt 누적 기록 (edge-trigger)",
        ),
        (
            "Open DVR",
            "RCS Tool List 창에서 해당 EQP_ID 의 Tool DVR(CCTV) 창을 자동 오픈",
        ),
        (
            "Zoom CH4",
            "DVR Player 창에서 Channel 4 를 확대하여 챔버 내부 영상만 단독 표시",
        ),
        (
            "Capture",
            "최대 8분(480초)간 100ms 간격으로 JPEG 프레임 저장 + summary.json/timeline.txt",
        ),
    ]

    step_top = Inches(1.5)
    step_height = Inches(1.55)
    step_width = Inches(2.45)
    gap = Inches(0.07)
    left = Inches(0.4)
    for idx, (title, detail) in enumerate(steps, start=1):
        _add_pipeline_step(
            slide,
            left=left,
            top=step_top,
            width=step_width,
            height=step_height,
            index=idx,
            title=title,
            detail=detail,
        )
        left += step_width + gap

    _add_section_card(
        slide,
        left=Inches(0.4),
        top=Inches(3.3),
        width=Inches(6.25),
        height=Inches(3.55),
        heading="핵심 모듈",
        bullets=[
            "align_fail_alarm.py — 감지 + 팝업 알림 전용 (1분 주기, 자동화 없음)",
            "monitor_align_fail.py — 감지 + RCS/CCTV/캡처 풀 자동화 (2분 주기)",
            "office_align_fail_alarm.py — 사내 알람 API 호출 & ALID=9006 필터",
            "workflow_select_tool_cctv.py — EQP_ID 로 Tool DVR 창 진입",
            "workflow_select_ch4_cctv.py — DVR Player 에서 Channel 4 확대",
            "capture_window_frames_ch4.py — 100ms 간격 JPEG 캡처 + 메타데이터 저장",
        ],
    )

    _add_section_card(
        slide,
        left=Inches(6.85),
        top=Inches(3.3),
        width=Inches(6.1),
        height=Inches(3.55),
        heading="운영 / 출력물",
        bullets=[
            "실행: uv run python poc/workflow_1/monitor_align_fail.py",
            "전제: 엔지니어가 RCS Tool List 창을 미리 열어둔 상태",
            "출력: 캡처 폴더당 ~4,800 JPEG 프레임 + summary.json + timeline.txt",
            "로그: poc/workflow_1/logs/align_fail_alarms.txt (감지 시각/EQP_ID/ALID)",
            "안전: 캡처 종료 후 DVR 창을 닫고 Tool List 창으로 포커스 복귀",
            "다음 단계: 수집된 프레임을 VLM(UI-Venus) 으로 Align Fail 원인 자동 분류",
        ],
    )

    _add_footer(
        slide,
        "Status: 사내 Windows 환경에서 동작 검증 중  •  다음 단계: VLM 기반 원인 자동 분석 연계",
    )


def _add_role_card(slide, *, left, top, width, height, role, model, responsibility, output, color):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    box.shadow.inherit = False

    role_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.35)
    )
    p = role_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = role
    _set_run(run, size=11, bold=True, color=color)

    model_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.45), width - Inches(0.4), Inches(0.35)
    )
    p = model_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = model
    _set_run(run, size=14, bold=True, color=NAVY)

    body_box = slide.shapes.add_textbox(
        left + Inches(0.2),
        top + Inches(0.85),
        width - Inches(0.4),
        height - Inches(1.0),
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_after = Pt(3)
    run = p.add_run()
    run.text = f"역할: {responsibility}"
    _set_run(run, size=10.5, color=TEXT)

    p = tf.add_paragraph()
    p.space_after = Pt(3)
    run = p.add_run()
    run.text = f"산출: {output}"
    _set_run(run, size=10.5, color=MUTED)


def build_slide_collab_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_title_band(
        slide,
        "여러 모델이 협력해서 자동화를 완성한다",
        "단일 모델의 한계를 역할 분담으로 보완 — UI 위치 / 정밀 클릭 / 텍스트 검증 / 정량 점수",
    )

    _add_section_card(
        slide,
        left=Inches(0.4),
        top=Inches(1.4),
        width=Inches(12.55),
        height=Inches(1.55),
        heading="왜 협력이 필요한가",
        bullets=[
            "단일 VLM 으로 '어디를 클릭할지' 를 한 번에 픽셀 정확도로 맞히기는 어렵다 —"
            " 전체 화면에서 좌표를 직접 찍으면 오차가 커서 실제 입력 자동화가 깨진다",
            "VLM 은 호출 간 기억이 없는(stateless) 모델이므로,"
            " '직전에 무엇을 했는가' 를 외부에서 보존하고 다음 모델로 넘겨야 한다",
            "그래서 한 모델이 모든 일을 하는 대신,"
            " 각 단계의 강점이 다른 모델을 파이프라인으로 엮어 GUI 자동화를 안정화한다",
        ],
    )

    role_top = Inches(3.1)
    role_height = Inches(2.55)
    role_width = Inches(3.05)
    gap = Inches(0.13)
    left = Inches(0.4)

    roles = [
        (
            "1. Coarse Locator",
            "UI-Venus 1.5 (8B)",
            "전체 화면에서 타겟 UI 요소의 대략적 bbox 를 찾아낸다",
            "bbox_1000 좌표 + 타겟 영역 crop 후보",
            NAVY,
        ),
        (
            "2. Fine Refiner",
            "MAI-UI",
            "Coarse bbox 주변을 확대한 zoom crop 에서 정밀 클릭 좌표를 찍는다",
            "(x, y) 픽셀 단위 클릭 포인트",
            ACCENT,
        ),
        (
            "3. Text Verifier",
            "PaddleOCR-VL 1.5",
            "입력 후 화면의 텍스트를 OCR 로 읽어 의도한 값이 들어갔는지 검증한다",
            "필드별 OCR 결과 + 일치 여부",
            NAVY,
        ),
        (
            "4. CV Scorer",
            "OpenCV (Template / SSIM / ORB)",
            "프레임-템플릿 유사도를 정량 점수화하여 stateless VLM 의 기억을 대신한다",
            "수치 매칭 점수 + best-frame 좌표",
            ACCENT,
        ),
    ]

    for role, model, responsibility, output, color in roles:
        _add_role_card(
            slide,
            left=left,
            top=role_top,
            width=role_width,
            height=role_height,
            role=role,
            model=model,
            responsibility=responsibility,
            output=output,
            color=color,
        )
        left += role_width + gap

    _add_section_card(
        slide,
        left=Inches(0.4),
        top=Inches(5.85),
        width=Inches(12.55),
        height=Inches(1.1),
        heading="협력으로 얻는 것",
        bullets=[
            "정확도: Coarse → Fine 분리로 클릭 좌표 오차를 ‘대략 영역’ 수준에서 ‘픽셀’ 수준으로 축소",
            "검증성: 액션 → OCR 검증 → 다음 액션 의 폐루프(closed-loop) 자동화로 무인 신뢰도 확보",
        ],
    )

    _add_footer(
        slide,
        "UI-Venus + MAI-UI + PaddleOCR-VL + OpenCV  =  GUI 자동화에 필요한 능력의 모듈식 조합",
    )


def build_slide_collab_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_title_band(
        slide,
        "협력 파이프라인: Coarse → Zoom → Fine → Verify",
        "RCS 로그인 자동화에서 실제 동작하는 다중 모델 협업 흐름",
    )

    steps = [
        (
            "Capture",
            "RCS 로그인 창을 캡처하여 WebP 로 인코딩 후 Flask 프록시로 전달",
        ),
        (
            "UI-Venus",
            "전체 이미지에서 'User ID 입력 필드' 의 coarse bbox(gold) 를 추출",
        ),
        (
            "Crop & Zoom",
            "Coarse bbox 주변에 padding 을 더해 잘라내고 확대하여 디테일 보존",
        ),
        (
            "MAI-UI",
            "Zoom crop 에서 refined 클릭 포인트(deepskyblue) 를 픽셀 단위로 결정",
        ),
        (
            "PaddleOCR-VL",
            "클릭 → 키 입력 후 화면을 다시 OCR 하여 의도한 값이 입력됐는지 검증",
        ),
    ]

    step_top = Inches(1.5)
    step_height = Inches(1.55)
    step_width = Inches(2.45)
    gap = Inches(0.07)
    left = Inches(0.4)
    for idx, (title, detail) in enumerate(steps, start=1):
        _add_pipeline_step(
            slide,
            left=left,
            top=step_top,
            width=step_width,
            height=step_height,
            index=idx,
            title=title,
            detail=detail,
        )
        left += step_width + gap

    _add_section_card(
        slide,
        left=Inches(0.4),
        top=Inches(3.3),
        width=Inches(6.25),
        height=Inches(3.55),
        heading="모델 간 데이터 흐름",
        bullets=[
            "UI-Venus 출력(bbox_1000) → 픽셀 bbox 환산 → crop box 계산 입력으로 전달",
            "Crop 이미지 + 동일한 target description → MAI-UI 입력으로 전달 (역할 키 재사용)",
            "MAI-UI 의 zoom 좌표 → crop 오프셋을 더해 원본 이미지 좌표계로 역변환",
            "원본 좌표계의 클릭 포인트 → pynput 으로 실제 클릭/키 입력 실행",
            "입력 후 캡처 → PaddleOCR-VL 로 텍스트 검증 → 다음 단계 진행 또는 재시도",
        ],
    )

    _add_section_card(
        slide,
        left=Inches(6.85),
        top=Inches(3.3),
        width=Inches(6.1),
        height=Inches(3.55),
        heading="구현 위치 / 핵심 파일",
        bullets=[
            "ui_venus_mai_locator.py — analyze_window_target() 가 2단계 파이프라인 오케스트레이션",
            "_run_ui_venus_coarse_bbox() / _run_mai_ui_refinement() — 각 모델 호출 분리",
            "_build_crop_box() — coarse bbox + padding 비율로 zoom 영역 산출",
            "prompts/prompt_login_rcs_ui_venus.py / prompt_login_rcs_mai_ui.py — 모델별 프롬프트 분리",
            "prompts/prompt_ocr_assist.py — PaddleOCR-VL 기반 텍스트 검증 프롬프트",
            "flask_vlm.py — service slug 별로 ui-venus / mai-ui / paddleocr-vl-1.5 라우팅",
        ],
    )

    _add_footer(
        slide,
        "Coarse Locator → Fine Refiner → Text Verifier  •  서로 다른 모델이 자기 강점만 책임진다",
    )


def _add_loop_node(slide, *, left, top, width, height, title, detail, color):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    box.shadow.inherit = False

    title_box = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), Inches(0.35)
    )
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, size=12, bold=True, color=color)

    body_box = slide.shapes.add_textbox(
        left + Inches(0.15),
        top + Inches(0.5),
        width - Inches(0.3),
        height - Inches(0.6),
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = detail
    _set_run(run, size=10, color=MUTED)


def _add_arrow(slide, *, left, top, width, height, color, direction="right"):
    shape_map = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
    }
    arrow = slide.shapes.add_shape(shape_map[direction], left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()


def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_title_band(
        slide,
        "다음 난관: 움직이는 SEM 화면에서 Align Key 자동 추적",
        "VLM 의 무기억(stateless) 한계를 이미지 프로세싱 스코어링으로 보완하는 탐색 루프",
    )

    _add_section_card(
        slide,
        left=Inches(0.4),
        top=Inches(1.4),
        width=Inches(6.25),
        height=Inches(2.0),
        heading="문제",
        bullets=[
            "SEM Monitor 화면은 스테이지 이동에 따라 실시간으로 바뀜 (정지 이미지가 아님)",
            "현재 화면이 Recipe 상의 Align Key 와 동일한지를 매 프레임 판단해야 함",
            "VLM 은 호출 간 기억이 없어 '이전에 본 위치/유사도' 를 스스로 추적하지 못함",
            "단일 VLM 호출만으로는 정량적 매칭 점수를 안정적으로 얻기 어려움",
        ],
    )

    _add_section_card(
        slide,
        left=Inches(6.85),
        top=Inches(1.4),
        width=Inches(6.1),
        height=Inches(2.0),
        heading="해결 전략",
        bullets=[
            "역할 분리: Image Processing = '얼마나 닮았나' (객관 점수),"
            " VLM = '다음에 어디를 볼까' (탐색 정책)",
            "매 프레임 Recipe Align Key vs SEM 화면을 점수화 →"
            " 외부 상태로 보존하여 VLM 의 무기억 보완",
            "점수가 임계값을 넘을 때까지 VLM 이 스테이지 이동 방향을 제안하며 루프",
            "Best-score 프레임/좌표를 저장하여 실패 시 fallback 으로 활용",
        ],
    )

    # 루프 다이어그램 (좌→우 4단계 + 하단 피드백 화살표)
    loop_top = Inches(3.7)
    loop_height = Inches(1.55)
    node_width = Inches(2.55)
    gap = Inches(0.35)
    arrow_width = Inches(0.3)
    arrow_height = Inches(0.45)

    nodes = [
        (
            "1. Capture",
            "SEM Monitor 창에서 현재 프레임 캡처 (CH4 캡처 파이프라인 재사용)",
            NAVY,
        ),
        (
            "2. Score (CV)",
            "Recipe Align Key 템플릿 vs 현재 프레임 → 유사도 점수\n"
            "(Template Match / SSIM / ORB feature)",
            ACCENT,
        ),
        (
            "3. Decide (VLM)",
            "점수 + 현재 프레임을 VLM 에 전달 → 일치 여부 판정 +\n"
            "다음 이동 방향(↑↓←→/배율) 제안",
            NAVY,
        ),
        (
            "4. Move Stage",
            "VLM 제안에 따라 SEM 스테이지 이동 명령 →\n"
            "임계값 도달까지 1단계로 회귀",
            ACCENT,
        ),
    ]

    left = Inches(0.4)
    node_lefts = []
    for idx, (title, detail, color) in enumerate(nodes):
        _add_loop_node(
            slide,
            left=left,
            top=loop_top,
            width=node_width,
            height=loop_height,
            title=title,
            detail=detail,
            color=color,
        )
        node_lefts.append(left)
        if idx < len(nodes) - 1:
            arrow_left = left + node_width + Inches(0.025)
            arrow_top = loop_top + (loop_height - arrow_height) / 2
            _add_arrow(
                slide,
                left=arrow_left,
                top=arrow_top,
                width=arrow_width,
                height=arrow_height,
                color=NAVY,
                direction="right",
            )
        left += node_width + gap

    # 피드백 루프 라벨 (하단)
    feedback_box = slide.shapes.add_textbox(
        Inches(0.4), loop_top + loop_height + Inches(0.1), Inches(12.55), Inches(0.35)
    )
    p = feedback_box.text_frame.paragraphs[0]
    p.alignment = 2  # center
    run = p.add_run()
    run.text = "↻  점수 < 임계값이면 1단계로 회귀 (Best-score 갱신하며 수렴까지 반복)"
    _set_run(run, size=11, bold=True, color=ACCENT)

    _add_section_card(
        slide,
        left=Inches(0.4),
        top=Inches(5.85),
        width=Inches(12.55),
        height=Inches(1.1),
        heading="기술 스택 & 검증 지표",
        bullets=[
            "이미지 매칭: OpenCV Template Matching · SSIM · ORB/AKAZE feature score 후보 비교",
            "수렴 지표: 매칭 점수 임계값 도달까지의 평균 루프 횟수 / 실패율 / Best-score fallback 정확도",
        ],
    )

    _add_footer(
        slide,
        "Stateless VLM + Stateful CV Score = Convergent Search Loop",
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_collab_overview(prs)
    build_slide_collab_pipeline(prs)
    build_slide_3(prs)

    prs.save(OUTPUT_PATH)
    print(f"[INFO] PPTX saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
