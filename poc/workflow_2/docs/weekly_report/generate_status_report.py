"""Workflow 2 매니저 보고용 진행 상황 리포트 생성기.

단계 데이터(WF1 압축 트랙 + WF2 본 트랙)를 정의하고,
같은 데이터에서 HTML과 PPTX 두 포맷으로 렌더링한다.

산출물:
    - workflow_2_status_report.html (외부 의존성 없는 단일 HTML)
    - workflow_2_status_report.pptx (16:9, 4슬라이드)

사용:
    uv run python poc/workflow_2/docs/weekly_report/generate_status_report.py
"""

from dataclasses import dataclass
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu


DOCS_DIR = Path(__file__).resolve().parent
HTML_PATH = DOCS_DIR / "workflow_2_status_report.html"
PPTX_PATH = DOCS_DIR / "workflow_2_status_report.pptx"
REPORT_DATE = date(2026, 6, 4)


STATUS_DONE = "done"
STATUS_IN_PROGRESS = "in_progress"
STATUS_PENDING = "pending"


STATUS_META = {
    STATUS_DONE: {
        "label": "완료",
        "icon": "✅",
        "border": "#2E7D32",
        "fill": "#E8F5E9",
        "text": "#1B5E20",
    },
    STATUS_IN_PROGRESS: {
        "label": "진행 중",
        "icon": "🔄",
        "border": "#ED6C02",
        "fill": "#FFF3E0",
        "text": "#B25500",
    },
    STATUS_PENDING: {
        "label": "대기",
        "icon": "⏳",
        "border": "#757575",
        "fill": "#F5F5F5",
        "text": "#424242",
    },
}


@dataclass
class Step:
    name: str
    status: str
    description: str
    current: bool = False


WF1_STEPS: list[Step] = [
    Step("RCS 로그인", STATUS_DONE, "VLM 기반 로그인 다이얼로그 자동 입력"),
    Step("Tool 선택", STATUS_DONE, "리스트 탭에서 측정 Tool 선택"),
    Step("Window/Frame Capture", STATUS_DONE, "RCS·CH4 윈도우 캡처 파이프라인"),
    Step("Popup/Cursor 감지", STATUS_DONE, "VLM으로 팝업·커서 상태 인식"),
    Step("Align Fail Alarm", STATUS_DONE, "정렬 실패 모니터링 + 알람"),
]


WF2_STEPS: list[Step] = [
    Step(
        "SEM 화면 영역 검출",
        STATUS_DONE,
        "모든 탐색의 기준점인 live SEM 박스를 검출, 장비별 기준 위치 저장",
    ),
    Step(
        "OK 버튼 검출·클릭",
        STATUS_DONE,
        "보정을 확정하는 마지막 동작, OCR 재확인으로 오클릭 차단",
    ),
    Step(
        "흰 박스 제거 (template)",
        STATUS_DONE,
        "고정 디지털값(히스토그램 섬) photometric 검출로 박스 안쪽만 crop, 최악 샘플 1/6→6/6 (합성검증)",
    ),
    Step(
        "십자(crosshair) 제거 (frame)",
        STATUS_DONE,
        "십자 검출→위치를 정답으로 기록→inpaint, 재검출+edge밀도로 제거 이중 검증 (합성검증)",
    ),
    Step(
        "Align Point 정확도 검증",
        STATUS_IN_PROGRESS,
        "2×2 비교(박스·십자 정리 효과 분리) + S전용 위치추정 검증(1발 명중률), 오피스 데이터 대기",
        current=True,
    ),
    Step(
        "실장비 FOV Search Loop",
        STATUS_PENDING,
        "검증 통과 후 실 SEM 연결 어댑터·안전장치 구현 → 매처 결과로 stage 이동",
    ),
]


NEXT_STEPS = [
    "(오피스) 골든셋 수집 — 성공(S) 측정 이미지를 recipe당 8~10장, 웨이퍼/lot/시간 분산",
    "정답 위생 먼저 확인 — 십자 검출률이 낮으면 그 위 수치 신뢰 불가, 검출기부터 점검",
    "본 판정 — golden_localization_eval.py 로 NEW(박스+inpaint)의 1발 명중률을 OLD 대비 측정, 통과 시 정식 채택",
    "임계 캘리브레이션 — success_vs_fail_compare.py 로 거짓양성 0 검증 후 fail 데이터 적용",
    "잔여 실패율로 VLM+CV 백업(gated escalation) 필요량 판단 → 필요시에만 구축",
]


# ──────────────────────────────────────────────────────────────────────
# HTML 렌더러
# ──────────────────────────────────────────────────────────────────────


def _wf1_chip_html(step: Step) -> str:
    return (
        '<div class="wf1-chip">'
        f'<span class="wf1-check">✓</span>'
        f'<span class="wf1-name">{step.name}</span>'
        "</div>"
    )


def _wf2_card_html(step: Step) -> str:
    meta = STATUS_META[step.status]
    pin = '<div class="pin">📍 현재</div>' if step.current else ""
    return (
        '<div class="wf2-card" '
        f'style="border-color:{meta["border"]};background:{meta["fill"]};color:{meta["text"]};">'
        f"{pin}"
        f'<div class="wf2-status">{meta["icon"]} {meta["label"]}</div>'
        f'<div class="wf2-name">{step.name}</div>'
        f'<div class="wf2-desc">{step.description}</div>'
        "</div>"
    )


def _legend_html() -> str:
    items = []
    for key in (STATUS_DONE, STATUS_IN_PROGRESS, STATUS_PENDING):
        meta = STATUS_META[key]
        items.append(
            '<span class="legend-item">'
            f'<span class="legend-swatch" style="background:{meta["fill"]};border-color:{meta["border"]};"></span>'
            f'{meta["icon"]} {meta["label"]}'
            "</span>"
        )
    return '<div class="legend">' + "".join(items) + "</div>"


def _status_table_html() -> str:
    rows = []
    for idx, step in enumerate(WF2_STEPS, start=1):
        meta = STATUS_META[step.status]
        marker = '<span class="row-pin">📍</span>' if step.current else ""
        rows.append(
            "<tr>"
            f'<td class="num">{idx}</td>'
            f'<td class="name">{step.name}{marker}</td>'
            f'<td class="status"><span class="badge" '
            f'style="background:{meta["fill"]};border-color:{meta["border"]};color:{meta["text"]};">'
            f'{meta["icon"]} {meta["label"]}</span></td>'
            f'<td class="desc">{step.description}</td>'
            "</tr>"
        )
    return (
        '<table class="status-table">'
        "<thead><tr><th>#</th><th>단계</th><th>상태</th><th>설명</th></tr></thead>"
        f'<tbody>{"".join(rows)}</tbody>'
        "</table>"
    )


def build_html() -> str:
    wf1_chips = '<div class="wf1-track">' + "".join(_wf1_chip_html(s) for s in WF1_STEPS) + "</div>"

    wf2_cards = []
    for i, step in enumerate(WF2_STEPS):
        wf2_cards.append(_wf2_card_html(step))
        if i < len(WF2_STEPS) - 1:
            wf2_cards.append('<div class="wf2-arrow">▶</div>')
    wf2_flow = '<div class="wf2-track">' + "".join(wf2_cards) + "</div>"

    next_items = "".join(f"<li>{item}</li>" for item in NEXT_STEPS)

    return f"""<!doctype html>
<html lang="ko">
<head>
<meta charset="utf-8">
<title>Workflow 2 진행 상황 보고</title>
<style>
  * {{ box-sizing: border-box; }}
  body {{
    font-family: 'Malgun Gothic', 'Apple SD Gothic Neo', 'Noto Sans KR', sans-serif;
    margin: 0; padding: 32px 48px; color: #212121; background: #FAFAFA;
    line-height: 1.55;
  }}
  h1 {{ font-size: 26px; margin: 0 0 4px 0; }}
  h2 {{ font-size: 18px; margin: 28px 0 12px 0; color: #455A64; border-bottom: 2px solid #CFD8DC; padding-bottom: 4px; }}
  .meta {{ color: #607D8B; font-size: 13px; margin-bottom: 12px; }}
  .summary {{ background: #fff; border-left: 4px solid #1976D2; padding: 12px 16px; margin: 8px 0 16px 0; font-size: 14px; }}

  /* WF1 압축 트랙 */
  .wf1-label {{ font-size: 12px; color: #757575; margin-bottom: 6px; }}
  .wf1-track {{
    display: flex; flex-wrap: wrap; gap: 8px; padding: 10px 14px;
    background: #ECEFF1; border-radius: 8px; align-items: center;
  }}
  .wf1-chip {{
    display: inline-flex; align-items: center; gap: 6px;
    padding: 4px 10px; background: #fff; border: 1px solid #B0BEC5;
    border-radius: 14px; font-size: 12px; color: #455A64;
  }}
  .wf1-check {{ color: #2E7D32; font-weight: bold; }}
  .wf1-name {{ font-weight: 500; }}

  /* WF2 본 트랙 */
  .wf2-track {{
    display: flex; align-items: stretch; gap: 6px; flex-wrap: nowrap;
    overflow-x: auto; padding: 18px 4px 8px 4px;
  }}
  .wf2-card {{
    flex: 1 1 0; min-width: 160px; position: relative;
    border: 2px solid; border-radius: 10px; padding: 14px 12px;
    background: #fff;
  }}
  .wf2-status {{ font-size: 12px; font-weight: 600; margin-bottom: 6px; opacity: 0.85; }}
  .wf2-name {{ font-size: 14px; font-weight: 700; margin-bottom: 6px; line-height: 1.35; }}
  .wf2-desc {{ font-size: 12px; line-height: 1.4; opacity: 0.85; }}
  .pin {{
    position: absolute; top: -14px; left: 50%; transform: translateX(-50%);
    background: #1565C0; color: #fff; font-size: 11px; font-weight: 700;
    padding: 2px 10px; border-radius: 12px; white-space: nowrap;
    box-shadow: 0 2px 6px rgba(0,0,0,0.15);
  }}
  .wf2-arrow {{
    display: flex; align-items: center; color: #90A4AE;
    font-size: 18px; padding: 0 2px;
  }}

  /* Legend */
  .legend {{ display: flex; gap: 16px; font-size: 12px; color: #455A64; margin-top: 8px; }}
  .legend-swatch {{
    display: inline-block; width: 14px; height: 14px; border: 1px solid;
    border-radius: 3px; margin-right: 4px; vertical-align: middle;
  }}

  /* Status table */
  .status-table {{
    width: 100%; border-collapse: collapse; background: #fff;
    font-size: 13px;
  }}
  .status-table th, .status-table td {{
    border: 1px solid #E0E0E0; padding: 8px 10px; text-align: left; vertical-align: top;
  }}
  .status-table th {{ background: #ECEFF1; font-weight: 600; color: #37474F; }}
  .status-table td.num {{ width: 36px; text-align: center; color: #757575; }}
  .status-table td.name {{ width: 220px; font-weight: 600; }}
  .status-table td.status {{ width: 110px; }}
  .badge {{
    display: inline-block; padding: 2px 8px; border: 1px solid; border-radius: 12px;
    font-size: 12px; font-weight: 600;
  }}
  .row-pin {{ margin-left: 6px; }}

  /* Next steps */
  .next ol {{ margin: 0; padding-left: 20px; }}
  .next li {{ margin-bottom: 6px; font-size: 13px; }}

  .footer {{ margin-top: 28px; font-size: 11px; color: #90A4AE; text-align: right; }}
</style>
</head>
<body>

<h1>Workflow 2 진행 상황 보고</h1>
<div class="meta">작성일: {REPORT_DATE.isoformat()} · 범위: poc/workflow_2/</div>

<div class="summary">
  Workflow 1 (RCS 로그인 → Align Fail Alarm)은 완성 단계에 있습니다.
  이어지는 <b>Workflow 2</b>는 현재 핵심 난제인 <b>Align Point 보정 정확도</b>에 집중하고 있습니다.
  매처를 교란하는 두 표식 — <b>등록 이미지의 흰 박스</b>·<b>측정 이미지의 십자</b> — 제거 방법을 만들었고
  (합성 검증 완료), <b>"이 정리가 실제로 좌표 명중률을 올리는가"</b>를 같은 잣대로 재는 검증 단계에
  진입했습니다. 실데이터(오피스) 판정이 다음 단계입니다.
</div>

<h2>① Process Flow</h2>
<div class="wf1-label">Workflow 1 · 완성 단계 (별도 보고)</div>
{wf1_chips}

<div style="text-align:center; color:#90A4AE; margin: 10px 0; font-size: 16px;">▼</div>

<div class="wf1-label">Workflow 2 (현재 작업 구간)</div>
{wf2_flow}
{_legend_html()}

<h2>② 진행 상황</h2>
{_status_table_html()}

<h2>③ Next Steps</h2>
<div class="next">
  <ol>{next_items}</ol>
</div>

<div class="footer">생성 스크립트: poc/workflow_2/docs/weekly_report/generate_status_report.py</div>

</body>
</html>
"""


# ──────────────────────────────────────────────────────────────────────
# PPTX 렌더러
# ──────────────────────────────────────────────────────────────────────


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_font(run, size=12, bold=False, color="#212121", font_name="Malgun Gothic"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _hex_to_rgb(color)
    run.font.name = font_name


def _add_text_box(slide, left, top, width, height, text, *, size=12, bold=False,
                  color="#212121", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color)
    return tb


def _add_step_shape(slide, left, top, width, height, step: Step, *, compact=False):
    meta = STATUS_META[step.status]
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(meta["fill"])
    shape.line.color.rgb = _hex_to_rgb(meta["border"])
    shape.line.width = Pt(2)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = f"{meta['icon']} {meta['label']}"
    _set_font(r1, size=9, bold=True, color=meta["text"])

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = step.name
    _set_font(r2, size=11 if not compact else 10, bold=True, color=meta["text"])

    if not compact:
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = step.description
        _set_font(r3, size=8, color=meta["text"])

    if step.current:
        pin_w = Inches(1.0)
        pin_h = Inches(0.28)
        pin_left = left + (width - pin_w) // 2
        pin_top = top - pin_h // 2
        pin = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pin_left, pin_top, pin_w, pin_h)
        pin.fill.solid()
        pin.fill.fore_color.rgb = _hex_to_rgb("#1565C0")
        pin.line.fill.background()
        pin.shadow.inherit = False
        ptf = pin.text_frame
        ptf.margin_left = Emu(0)
        ptf.margin_right = Emu(0)
        ptf.margin_top = Emu(0)
        ptf.margin_bottom = Emu(0)
        pp = ptf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = "📍 현재"
        _set_font(pr, size=9, bold=True, color="#FFFFFF")


def _add_arrow(slide, x1, y1, x2, y2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = _hex_to_rgb("#90A4AE")
    line.line.width = Pt(1.5)


def _slide_title(slide, title: str, subtitle: str = ""):
    _add_text_box(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.5),
                  title, size=22, bold=True, color="#212121")
    if subtitle:
        _add_text_box(slide, Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.3),
                      subtitle, size=11, color="#607D8B")


def build_pptx() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1: Title ──────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    _add_text_box(s1, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.0),
                  "Workflow 2 진행 상황 보고", size=40, bold=True, color="#1A237E",
                  align=PP_ALIGN.CENTER)
    _add_text_box(s1, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.5),
                  "Align Point 보정 정확도 — 표식 제거 + 검증", size=20, color="#37474F",
                  align=PP_ALIGN.CENTER)
    _add_text_box(s1, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.4),
                  f"작성일 {REPORT_DATE.isoformat()}  ·  범위 poc/workflow_2/",
                  size=14, color="#78909C", align=PP_ALIGN.CENTER)
    _add_text_box(s1, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.4),
                  "Workflow 1 (RCS 로그인 ~ Align Fail Alarm)은 완성 단계에 있으며, 본 보고는 Workflow 2의 진행 현황입니다",
                  size=12, color="#90A4AE", align=PP_ALIGN.CENTER)

    # ── Slide 2: Process Flow ────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    _slide_title(s2, "① Process Flow", "Workflow 1 (완성 단계) → Workflow 2 (현재 작업)")

    # WF1 압축 트랙
    _add_text_box(s2, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.3),
                  "Workflow 1 · 완성 단계 (별도 보고)", size=11, color="#757575")

    wf1_left = Inches(0.5)
    wf1_top = Inches(1.6)
    wf1_total_w = Inches(12.3)
    wf1_card_w = Emu((wf1_total_w - Inches(0.4)) // len(WF1_STEPS))
    wf1_card_h = Inches(0.55)
    wf1_gap = Inches(0.1)
    actual_card_w = Emu((int(wf1_total_w) - int(wf1_gap) * (len(WF1_STEPS) - 1)) // len(WF1_STEPS))
    for i, step in enumerate(WF1_STEPS):
        left = Emu(int(wf1_left) + i * (int(actual_card_w) + int(wf1_gap)))
        chip = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, wf1_top, actual_card_w, wf1_card_h)
        chip.fill.solid()
        chip.fill.fore_color.rgb = _hex_to_rgb("#ECEFF1")
        chip.line.color.rgb = _hex_to_rgb("#B0BEC5")
        chip.line.width = Pt(1)
        chip.shadow.inherit = False
        tf = chip.text_frame
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"✓ {step.name}"
        _set_font(r, size=10, bold=True, color="#455A64")

    # 아래 화살표
    arrow_x = Inches(6.65)
    _add_arrow(s2, arrow_x, Inches(2.3), arrow_x, Inches(2.75))
    _add_text_box(s2, Inches(6.0), Inches(2.45), Inches(1.5), Inches(0.3),
                  "▼", size=14, color="#90A4AE", align=PP_ALIGN.CENTER)

    # WF2 본 트랙
    _add_text_box(s2, Inches(0.5), Inches(2.95), Inches(12.3), Inches(0.3),
                  "Workflow 2 · 현재 작업 구간", size=11, bold=True, color="#1565C0")

    wf2_left = Inches(0.5)
    wf2_top = Inches(3.6)
    wf2_total_w = Inches(12.3)
    wf2_card_h = Inches(2.4)
    wf2_gap = Inches(0.15)
    wf2_card_w = Emu((int(wf2_total_w) - int(wf2_gap) * (len(WF2_STEPS) - 1)) // len(WF2_STEPS))
    for i, step in enumerate(WF2_STEPS):
        left = Emu(int(wf2_left) + i * (int(wf2_card_w) + int(wf2_gap)))
        _add_step_shape(s2, left, wf2_top, wf2_card_w, wf2_card_h, step)
        if i < len(WF2_STEPS) - 1:
            ax = Emu(int(left) + int(wf2_card_w))
            ay = Emu(int(wf2_top) + int(wf2_card_h) // 2)
            _add_arrow(s2, ax, ay, Emu(int(ax) + int(wf2_gap)), ay)

    # 범례
    legend_top = Inches(6.3)
    legend_left = Inches(0.5)
    for i, key in enumerate((STATUS_DONE, STATUS_IN_PROGRESS, STATUS_PENDING)):
        meta = STATUS_META[key]
        sw_left = Emu(int(legend_left) + i * int(Inches(2.0)))
        sw = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, sw_left, legend_top, Inches(0.2), Inches(0.2))
        sw.fill.solid()
        sw.fill.fore_color.rgb = _hex_to_rgb(meta["fill"])
        sw.line.color.rgb = _hex_to_rgb(meta["border"])
        sw.shadow.inherit = False
        _add_text_box(s2, Emu(int(sw_left) + int(Inches(0.3))), legend_top,
                      Inches(1.6), Inches(0.25),
                      f"{meta['icon']} {meta['label']}", size=10, color="#455A64")

    # ── Slide 3: 상태 표 ─────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    _slide_title(s3, "② 진행 상황", f"Workflow 2 {len(WF2_STEPS)}개 단계 상세")

    headers = ["#", "단계", "상태", "설명"]
    n_rows = len(WF2_STEPS) + 1
    n_cols = len(headers)
    table_left = Inches(0.5)
    table_top = Inches(1.3)
    table_w = Inches(12.3)
    table_h = Inches(5.4)
    table_shape = s3.shapes.add_table(n_rows, n_cols, table_left, table_top, table_w, table_h)
    table = table_shape.table
    table.columns[0].width = Inches(0.5)
    table.columns[1].width = Inches(3.3)
    table.columns[2].width = Inches(1.6)
    table.columns[3].width = Inches(6.9)

    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex_to_rgb("#37474F")
        tf = cell.text_frame
        tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = header
        _set_font(r, size=12, bold=True, color="#FFFFFF")

    for ri, step in enumerate(WF2_STEPS, start=1):
        meta = STATUS_META[step.status]
        row_bg = "#FFFFFF" if ri % 2 == 1 else "#FAFAFA"
        if step.current:
            row_bg = "#E3F2FD"

        values = [
            str(ri),
            step.name + ("  📍" if step.current else ""),
            f"{meta['icon']} {meta['label']}",
            step.description,
        ]
        aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT]
        bolds = [False, True, True, False]
        sizes = [11, 12, 11, 11]
        colors = ["#616161", "#212121", meta["text"], "#424242"]

        for ci in range(n_cols):
            cell = table.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _hex_to_rgb(row_bg)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = aligns[ci]
            r = p.add_run()
            r.text = values[ci]
            _set_font(r, size=sizes[ci], bold=bolds[ci], color=colors[ci])

    # ── Slide 4: Next Steps ─────────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    _slide_title(s4, "③ Next Steps", "단기 우선순위")

    tb = s4.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(NEXT_STEPS):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
        r = p.add_run()
        r.text = f"{i + 1}. {item}"
        _set_font(r, size=18, color="#263238")

    return prs


# ──────────────────────────────────────────────────────────────────────


def main() -> None:
    html = build_html()
    HTML_PATH.write_text(html, encoding="utf-8")
    print(f"[INFO] HTML 생성: {HTML_PATH}")

    prs = build_pptx()
    prs.save(str(PPTX_PATH))
    print(f"[INFO] PPTX 생성: {PPTX_PATH}")


if __name__ == "__main__":
    main()
